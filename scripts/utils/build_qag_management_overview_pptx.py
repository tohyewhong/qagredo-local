#!/usr/bin/env python3
"""One-slide management summary PPTX for QAG."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

_FILL = RGBColor(255, 255, 255)
_LINE = RGBColor(80, 90, 100)
_TEXT = RGBColor(30, 33, 40)
_EDGE = RGBColor(9, 105, 218)


def _box(slide, left, top, w, h, text: str, *, size: float = 14) -> object:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = _FILL
    sh.line.color.rgb = _LINE
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.name = "Calibri"
    p.font.color.rgb = _TEXT
    p.alignment = PP_ALIGN.CENTER
    return sh


def _dia(slide, left, top, side, text: str) -> object:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.FLOWCHART_DECISION, left, top, side, side
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = _FILL
    sh.line.color.rgb = _LINE
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.name = "Calibri"
    p.font.color.rgb = _TEXT
    p.alignment = PP_ALIGN.CENTER
    return sh


def _arrow(slide, x1, y1, x2, y2, label: str | None = None) -> None:
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, int(x1), int(y1), int(x2), int(y2)
    )
    c.line.color.rgb = _EDGE
    c.line.width = Pt(1.5)
    if label:
        mx, my = (int(x1) + int(x2)) // 2, (int(y1) + int(y2)) // 2
        tb = slide.shapes.add_textbox(
            Emu(mx - 400000), Emu(my - 90000), Emu(800000), Emu(180000)
        )
        tb.fill.background()
        tb.line.fill.background()
        p = tb.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.name = "Calibri"
        p.font.color.rgb = _EDGE


def build(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = int(13.3333 * 914400)
    prs.slide_height = int(7.5 * 914400)
    s = prs.slides.add_slide(prs.slide_layouts[6])

    t = s.shapes.add_textbox(
        Inches(0.4), Inches(0.35), Inches(12.5), Inches(0.9)
    )
    t.fill.background()
    t.line.fill.background()
    tp = t.text_frame.paragraphs[0]
    tp.text = "QAG — management view (simplified)"
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.name = "Calibri"
    tp.font.color.rgb = _TEXT

    sub = s.shapes.add_textbox(
        Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.55)
    )
    sub.fill.background()
    sub.line.fill.background()
    sp = sub.text_frame.paragraphs[0]
    sp.text = (
        "Per document: AI first drafts a list of questions, then answers "
        "each one. Quality check can trigger a replacement question for a "
        "slot (not only a new answer). Then structured JSON out. "
        "Two AI services + pipeline (typical offline layout)."
    )
    sp.font.size = Pt(14)
    sp.font.name = "Calibri"
    sp.font.color.rgb = RGBColor(90, 95, 105)

    y = Inches(2.15)
    w, h = Inches(1.62), Inches(0.95)
    g = Inches(0.22)
    x0 = Inches(0.35)

    def cx(i):
        return x0 + i * (w + g)

    a = _box(s, cx(0), y, w, h, "Documents\n(JSON / JSONL)", size=13)
    b = _box(s, cx(1), y, w, h, "Generate\nquestions\n(N per doc)", size=12)
    c = _box(s, cx(2), y, w, h, "Generate\nanswers", size=13)
    d = _box(s, cx(3), y, w, h, "Quality\ncheck", size=13)
    e = _dia(
        s,
        cx(4) + Inches(0.08),
        y + Inches(0.06),
        Inches(0.78),
        "OK?",
    )
    f = _box(
        s,
        cx(5) + Inches(0.15),
        y,
        Inches(1.55),
        h,
        "Save\nanalysis JSON",
        size=13,
    )
    gh = _box(
        s,
        cx(6) + Inches(0.25),
        y,
        Inches(1.45),
        h,
        "Publish to\nGitHub",
        size=13,
    )

    def mr(sh):
        return Emu(sh.left + sh.width), Emu(sh.top + sh.height // 2)

    def ml(sh):
        return Emu(sh.left), Emu(sh.top + sh.height // 2)

    def mb(sh):
        return Emu(sh.left + sh.width // 2), Emu(sh.top + sh.height)

    def mt(sh):
        return Emu(sh.left + sh.width // 2), Emu(sh.top)

    _arrow(s, *mr(a), *ml(b))
    _arrow(s, *mr(b), *ml(c))
    _arrow(s, *mr(c), *ml(d))
    _arrow(s, *mr(d), *ml(e))
    _arrow(s, *mr(e), *ml(f), "yes")
    _arrow(s, *mr(f), *ml(gh))

    # Not OK: replacement question for this slot, then answer again
    rq = _box(
        s,
        cx(2) - Inches(0.15),
        y + Inches(1.12),
        Inches(2.05),
        Inches(0.72),
        "If not OK:\nreplacement question\n(same slot)",
        size=11,
    )
    _arrow(
        s,
        Emu(e.left + e.width // 2),
        Emu(e.top + e.height),
        Emu(rq.left + rq.width // 2),
        Emu(rq.top),
    )
    elb = s.shapes.add_textbox(
        Inches(6.35), Inches(3.35), Inches(0.55), Inches(0.32)
    )
    elb.fill.background()
    elb.line.fill.background()
    elp = elb.text_frame.paragraphs[0]
    elp.text = "no"
    elp.font.size = Pt(9)
    elp.font.color.rgb = _EDGE
    elp.font.name = "Calibri"

    _arrow(
        s,
        Emu(rq.left + rq.width // 2),
        Emu(rq.top + rq.height),
        Emu(c.left + c.width // 2),
        Emu(c.top),
    )
    rlb = s.shapes.add_textbox(
        Inches(3.15), Inches(2.95), Inches(1.55), Inches(0.35)
    )
    rlb.fill.background()
    rlb.line.fill.background()
    rlp = rlb.text_frame.paragraphs[0]
    rlp.text = "then answer again"
    rlp.font.size = Pt(9)
    rlp.font.color.rgb = _EDGE
    rlp.font.name = "Calibri"

    note = s.shapes.add_textbox(
        Inches(0.4), Inches(4.25), Inches(12.5), Inches(2.65)
    )
    note.fill.background()
    note.line.fill.background()
    nf = note.text_frame
    nf.word_wrap = True
    p0 = nf.paragraphs[0]
    p0.text = "Talking points"
    p0.font.bold = True
    p0.font.size = Pt(16)
    p0.font.name = "Calibri"
    for line, sz in [
        (
            "• Questions: AI proposes N questions per document first; each "
            "slot is one question answered and checked.",
            13,
        ),
        (
            "• If weak: replacement question for that slot (then new answer), "
            "within configured limits—not only re-answering the same text.",
            13,
        ),
        (
            "• Output: timestamped folder, one analysis file per document; "
            "optional modes for training-friendly extracts.",
            13,
        ),
        (
            "• Trust: separate checker AI from writer AI; small embedding "
            "model for fast similarity.",
            13,
        ),
        (
            "• Detail: see QAG_Software_Engineering_Diagrams.html "
            "in /docs.",
            12,
        ),
    ]:
        p = nf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.name = "Calibri"
        p.space_before = Pt(6)

    foot = s.shapes.add_textbox(
        Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.4)
    )
    foot.fill.background()
    foot.line.fill.background()
    fp = foot.text_frame.paragraphs[0]
    fp.text = (
        "Regenerate: "
        "python3 scripts/utils/build_qag_management_overview_pptx.py"
    )
    fp.font.size = Pt(10)
    fp.font.color.rgb = _LINE

    prs.save(str(output))
    print(f"Wrote {output}")


def main() -> None:
    ap = argparse.ArgumentParser()
    root = Path(__file__).resolve().parents[2]
    d = root / "docs" / "QAG_Management_Overview.pptx"
    ap.add_argument("-o", "--output", type=Path, default=d)
    build(ap.parse_args().output.resolve())


if __name__ == "__main__":
    main()
