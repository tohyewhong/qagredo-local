#!/usr/bin/env python3
"""Build editable PPTX: QAG pipeline (2 slides — core + save)."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

_FILL = RGBColor(241, 245, 249)
_LINE = RGBColor(71, 85, 105)
_TEXT = RGBColor(15, 23, 42)
_EDGE = RGBColor(37, 99, 235)
_EDGE_NO = RGBColor(220, 38, 38)


def _term(slide, left, top, w, h, text: str) -> object:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.FLOWCHART_TERMINATOR, left, top, w, h
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = _FILL
    sh.line.color.rgb = _LINE
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.name = "Calibri"
    p.font.color.rgb = _TEXT
    p.alignment = PP_ALIGN.CENTER
    return sh


def _proc(slide, left, top, w, h, text: str) -> object:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.FLOWCHART_PROCESS, left, top, w, h
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = _FILL
    sh.line.color.rgb = _LINE
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(7.5)
    p.font.name = "Calibri"
    p.font.color.rgb = _TEXT
    p.alignment = PP_ALIGN.CENTER
    return sh


def _dec(slide, left, top, side, text: str) -> object:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.FLOWCHART_DECISION, left, top, side, side
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = _FILL
    sh.line.color.rgb = _LINE
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(6.5)
    p.font.name = "Calibri"
    p.font.color.rgb = _TEXT
    p.alignment = PP_ALIGN.CENTER
    return sh


def _arrow(
    slide,
    x1: Emu,
    y1: Emu,
    x2: Emu,
    y2: Emu,
    *,
    dashed: bool = False,
    label: str | None = None,
    color: RGBColor | None = None,
) -> None:
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, int(x1), int(y1), int(x2), int(y2)
    )
    conn.line.color.rgb = color or _EDGE
    conn.line.width = Pt(1.1)
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if label:
        mx = (int(x1) + int(x2)) // 2
        my = (int(y1) + int(y2)) // 2
        tb = slide.shapes.add_textbox(
            Emu(mx - 457200), Emu(my - 101600), Emu(914400), Emu(203200)
        )
        tb.fill.background()
        tb.line.fill.background()
        p = tb.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(7)
        p.font.name = "Calibri"
        p.font.color.rgb = color or _EDGE


def _curve(slide, x1: Emu, y1: Emu, x2: Emu, y2: Emu) -> None:
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.CURVE, int(x1), int(y1), int(x2), int(y2)
    )
    conn.line.color.rgb = _EDGE
    conn.line.width = Pt(1.25)


def _footnote(slide, text: str) -> None:
    foot = slide.shapes.add_textbox(
        Inches(0.35), Inches(7.02), Inches(12.6), Inches(0.42)
    )
    foot.fill.background()
    foot.line.fill.background()
    fp = foot.text_frame.paragraphs[0]
    fp.text = text
    fp.font.size = Pt(8)
    fp.font.name = "Calibri"
    fp.font.color.rgb = _LINE


def _build_slide_core(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.1), Inches(12.6), Inches(0.45)
    )
    title.fill.background()
    title.line.fill.background()
    tp = title.text_frame.paragraphs[0]
    tp.text = "QAG pipeline — core loop (per document)"
    tp.font.size = Pt(17)
    tp.font.bold = True
    tp.font.name = "Calibri"
    tp.font.color.rgb = _TEXT

    x_m = Inches(5.15)
    x_r = Inches(8.95)
    w_t, h_t = Inches(1.45), Inches(0.4)
    w_p, h_p = Inches(1.72), Inches(0.38)
    d_side = Inches(0.58)
    y0 = Inches(0.68)

    start = _term(slide, x_m + Inches(0.12), y0, w_t, h_t, "Start")
    load = _proc(
        slide, x_m, y0 + Inches(0.48), w_p, h_p, "Load JSONL\n& normalize",
    )
    d_doc = _dec(
        slide, x_m + Inches(0.05), y0 + Inches(0.95), d_side,
        "Another\ndocument?",
    )
    end_n = _term(
        slide, Inches(10.6), y0 + Inches(0.95), Inches(1.25), h_t,
        "Pipeline\nend",
    )
    gen_q = _proc(
        slide, x_m - Inches(0.05), y0 + Inches(1.58), Inches(1.82), h_p,
        "generate_questions\n(N per doc)",
    )
    d_slot = _dec(
        slide, x_m + Inches(0.05), y0 + Inches(2.08), d_side,
        "For each\nslot 1..N?",
    )
    gen_a = _proc(
        slide, x_m, y0 + Inches(2.68), w_p, h_p,
        "generate_answers\n(slot)",
    )
    rep_q = _proc(
        slide, x_r - Inches(0.05), y0 + Inches(2.68), Inches(1.48), h_p,
        "generate_questions\n(×1 replace)",
    )
    grade = _proc(
        slide, x_m + Inches(0.02), y0 + Inches(3.18), Inches(1.68), h_p,
        "grade + pair +\ngrounding check",
    )
    d_gate = _dec(
        slide, x_m + Inches(0.05), y0 + Inches(3.68), d_side,
        "Passes\nquality gate?",
    )
    d_rep = _dec(
        slide, x_r + Inches(0.02), y0 + Inches(3.68), d_side,
        "Regeneration\nrounds left?",
    )
    keep = _proc(
        slide, x_r + Inches(0.05), y0 + Inches(4.45), Inches(1.38),
        Inches(0.36), "Keep last\nattempt",
    )
    d_more = _dec(
        slide, x_m + Inches(0.05), y0 + Inches(4.95), d_side,
        "More\nslots?",
    )
    save = _proc(
        slide, x_m - Inches(0.02), y0 + Inches(5.55), Inches(1.78),
        Inches(0.4), "save_results +\nGitHub publish",
    )

    def mb(sh):
        return Emu(sh.left + sh.width // 2), Emu(sh.top + sh.height)

    def mt(sh):
        return Emu(sh.left + sh.width // 2), Emu(sh.top)

    _arrow(slide, *mb(start), *mt(load))
    _arrow(slide, *mb(load), *mt(d_doc))
    _arrow(
        slide,
        Emu(d_doc.left + d_doc.width),
        Emu(d_doc.top + d_doc.height // 2),
        Emu(end_n.left),
        Emu(end_n.top + end_n.height // 2),
        label="no",
    )
    _arrow(slide, *mb(d_doc), *mt(gen_q), label="yes")
    _arrow(slide, *mb(gen_q), *mt(d_slot))
    _arrow(slide, *mb(d_slot), *mt(gen_a), label="yes")
    _arrow(
        slide,
        Emu(d_slot.left + d_slot.width),
        Emu(d_slot.top + d_slot.height // 2),
        Emu(save.left + save.width // 2),
        Emu(save.top),
        label="no",
    )
    _arrow(slide, *mb(gen_a), *mt(grade))
    _arrow(slide, *mb(grade), *mt(d_gate))

    _arrow(
        slide,
        Emu(d_gate.left + d_gate.width // 2),
        Emu(d_gate.top + d_gate.height),
        Emu(d_more.left + d_more.width // 2),
        Emu(d_more.top),
        label="yes",
    )
    _arrow(
        slide,
        Emu(d_gate.left + d_gate.width),
        Emu(d_gate.top + d_gate.height // 2),
        Emu(d_rep.left),
        Emu(d_rep.top + d_rep.height // 2),
        label="no",
        color=_EDGE_NO,
    )

    _arrow(
        slide,
        Emu(d_rep.left + d_rep.width // 2),
        Emu(d_rep.top),
        Emu(rep_q.left + rep_q.width // 2),
        Emu(rep_q.top + rep_q.height),
        label="yes",
        color=_EDGE_NO,
    )
    _curve(
        slide,
        Emu(rep_q.left + rep_q.width // 2),
        Emu(rep_q.top),
        Emu(gen_a.left + gen_a.width // 2),
        Emu(gen_a.top),
    )
    tb = slide.shapes.add_textbox(
        Inches(7.35), Inches(2.35), Inches(0.85), Inches(0.25)
    )
    tb.fill.background()
    tb.line.fill.background()
    tp2 = tb.text_frame.paragraphs[0]
    tp2.text = "re-answer"
    tp2.font.size = Pt(6.5)
    tp2.font.name = "Calibri"
    tp2.font.color.rgb = _EDGE
    tp2.font.italic = True

    _arrow(
        slide, *mb(d_rep), *mt(keep), label="no", color=_EDGE_NO,
    )
    _arrow(
        slide,
        Emu(keep.left),
        Emu(keep.top + keep.height // 2),
        Emu(d_more.left + d_more.width // 2),
        Emu(d_more.top + d_more.height),
    )

    _curve(
        slide,
        Emu(d_more.left),
        Emu(d_more.top + d_more.height // 2),
        Emu(d_slot.left),
        Emu(d_slot.top + d_slot.height // 2),
    )
    ylb = slide.shapes.add_textbox(
        Inches(3.55), Inches(5.05), Inches(0.55), Inches(0.24)
    )
    ylb.fill.background()
    ylb.line.fill.background()
    yp = ylb.text_frame.paragraphs[0]
    yp.text = "yes"
    yp.font.size = Pt(7)
    yp.font.name = "Calibri"
    yp.font.color.rgb = _EDGE

    _arrow(slide, *mb(d_more), *mt(save), label="no")
    _curve(
        slide,
        Emu(save.left + save.width // 2),
        Emu(save.top),
        Emu(d_doc.left + d_doc.width // 2),
        Emu(d_doc.top + d_doc.height),
    )
    tloop = slide.shapes.add_textbox(
        Inches(5.0), Inches(1.05), Inches(1.2), Inches(0.22)
    )
    tloop.fill.background()
    tloop.line.fill.background()
    tlp = tloop.text_frame.paragraphs[0]
    tlp.text = "next doc"
    tlp.font.size = Pt(6.5)
    tlp.font.name = "Calibri"
    tlp.font.color.rgb = _EDGE

    _footnote(
        slide,
        "Slide 2: save / filter options. Regenerate: "
        "python3 scripts/utils/build_qag_pipeline_flowchart_pptx.py",
    )


def _build_slide_save(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.1), Inches(12.6), Inches(0.45)
    )
    title.fill.background()
    title.line.fill.background()
    tp = title.text_frame.paragraphs[0]
    tp.text = "QAG — save & output shaping (after slot loop)"
    tp.font.size = Pt(17)
    tp.font.bold = True
    tp.font.name = "Calibri"
    tp.font.color.rgb = _TEXT

    x0 = Inches(4.9)
    w_p, h_p = Inches(1.75), Inches(0.4)
    d_side = Inches(0.58)
    y0 = Inches(1.85)

    in_p = _proc(
        slide, x0, y0, w_p, h_p,
        "QA pairs ready\n(for one document)",
    )
    d_gr = _dec(
        slide, x0 + Inches(0.08), y0 + Inches(0.62), d_side,
        "save_grounded\nonly?",
    )
    filt = _proc(
        slide, Inches(2.85), y0 + Inches(1.25), Inches(1.38), h_p,
        "Filter pairs\nby gate",
    )
    pad = _proc(
        slide, Inches(7.55), y0 + Inches(1.25), Inches(1.38), h_p,
        "Pad\nplaceholders",
    )
    d_min = _dec(
        slide, x0 + Inches(0.08), y0 + Inches(2.15), d_side,
        "minimal_qa\noutput?",
    )
    save_min = _proc(
        slide, Inches(3.25), y0 + Inches(2.85), Inches(1.5), Inches(0.38),
        "save minimal JSON\n(+ GitHub)",
    )
    save_full = _proc(
        slide, Inches(6.75), y0 + Inches(2.85), Inches(1.5), Inches(0.38),
        "save full JSON\n(+ GitHub)",
    )

    def mb(sh):
        return Emu(sh.left + sh.width // 2), Emu(sh.top + sh.height)

    def mt(sh):
        return Emu(sh.left + sh.width // 2), Emu(sh.top)

    _arrow(slide, *mb(in_p), *mt(d_gr))
    _arrow(
        slide,
        Emu(d_gr.left),
        Emu(d_gr.top + d_gr.height // 2),
        Emu(filt.left + filt.width),
        Emu(filt.top + filt.height // 2),
        label="yes",
    )
    _arrow(
        slide,
        Emu(d_gr.left + d_gr.width),
        Emu(d_gr.top + d_gr.height // 2),
        Emu(pad.left),
        Emu(pad.top + pad.height // 2),
        label="no",
    )
    _arrow(slide, *mb(filt), *mt(d_min))
    _arrow(slide, *mb(pad), *mt(d_min))
    _arrow(
        slide,
        Emu(d_min.left),
        Emu(d_min.top + d_min.height // 2),
        Emu(save_min.left + save_min.width),
        Emu(save_min.top + save_min.height // 2),
        label="yes",
    )
    _arrow(
        slide,
        Emu(d_min.left + d_min.width),
        Emu(d_min.top + d_min.height // 2),
        Emu(save_full.left),
        Emu(save_full.top + save_full.height // 2),
        label="no",
    )

    note = slide.shapes.add_textbox(
        Inches(0.9), Inches(5.15), Inches(11.5), Inches(1.35)
    )
    note.fill.background()
    note.line.fill.background()
    np = note.text_frame.paragraphs[0]
    np.text = (
        "Optional filters before write — same run_qa_pipeline stage, split "
        "here so the main flowchart stays readable."
    )
    np.font.size = Pt(10)
    np.font.name = "Calibri"
    np.font.color.rgb = _LINE

    _footnote(
        slide,
        "Regenerate: python3 scripts/utils/"
        "build_qag_pipeline_flowchart_pptx.py",
    )


def build(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = int(13.3333 * 914400)
    prs.slide_height = int(7.5 * 914400)
    _build_slide_core(prs)
    _build_slide_save(prs)
    prs.save(str(output))
    print(f"Wrote {output}")


def main() -> None:
    ap = argparse.ArgumentParser(
        description="Build QAG pipeline flowchart PPTX (2 slides).",
    )
    root = Path(__file__).resolve().parents[2]
    default_out = (
        root / "docs" / "architecture" / "diagrams" /
        "QAG_Pipeline_Flowchart_editable.pptx"
    )
    ap.add_argument(
        "-o",
        "--output",
        type=Path,
        default=default_out,
        help="Output .pptx path",
    )
    args = ap.parse_args()
    build(args.output.resolve())


if __name__ == "__main__":
    main()
