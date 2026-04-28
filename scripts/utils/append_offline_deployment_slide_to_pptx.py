#!/usr/bin/env python3
"""Append an idempotent 'offline deployment / requirements sync' slide to doc PPTX files."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

MARKER = "Offline deployment — keep artifacts in sync"

BODY = """• requirements.txt in the bundle must match packages inside qagredo-v1.tar
• After requirements change: rebuild image, docker save; offline: docker rmi qagredo-v1:latest then docker load -i qagredo-v1.tar
• In qagredo_host/: bash verify_offline_deployment.sh after docker load
• make_qagredo_bundle.sh checks bundle vs local qagredo-v1:latest when the image exists
• Pipeline uses /usr/local/bin/python inside the container (see run.sh)

Authoritative guide: docs/OFFLINE_SETUP_GUIDE.md"""


def _last_slide_has_marker(prs: Presentation) -> bool:
    if not prs.slides:
        return False
    slide = prs.slides[-1]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            if MARKER in (p.text or ""):
                return True
    return False


def append_slide(path: Path) -> str:
    prs = Presentation(str(path))
    if _last_slide_has_marker(prs):
        return "skip (already updated)"
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    title.text = MARKER
    if body and body.has_text_frame:
        body.text_frame.text = BODY
        for p in body.text_frame.paragraphs:
            p.font.size = Pt(14)
    title.text_frame.paragraphs[0].font.size = Pt(24)
    prs.save(str(path))
    return "ok"


def main() -> int:
    repo = Path(__file__).resolve().parents[2]
    targets: list[Path] = []
    for sub in ("docs",):
        root = repo / sub
        if not root.is_dir():
            continue
        for p in root.rglob("*.pptx"):
            rp = str(p.resolve())
            if "/.venv/" in rp or "/.tmp/" in rp or "/tmp_validation/" in rp:
                continue
            if "/data/" in rp:
                continue
            targets.append(p)
    failed = 0
    for p in sorted(targets):
        try:
            status = append_slide(p)
            print(f"{status}: {p.relative_to(repo)}")
        except Exception as e:
            print(f"FAIL: {p.relative_to(repo)}: {e}", file=sys.stderr)
            failed += 1
    return 1 if failed else 0


if __name__ == "__main__":
    raise SystemExit(main())
