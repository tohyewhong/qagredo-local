#!/usr/bin/env python3
"""
Build editable PPTX for QAG per-slot flow.

Uses native shapes + connectors (editable in PowerPoint).
Includes LLM judge grading (legacy hybrid mode is LLM-only).
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

_FILL = RGBColor(241, 245, 249)
_LINE = RGBColor(100, 116, 139)
_TEXT = RGBColor(15, 23, 42)
_EDGE = RGBColor(71, 85, 105)


def _box(slide, left, top, w, h, text: str) -> object:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = _FILL
    sh.line.color.rgb = _LINE
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.name = "Calibri"
    p.font.color.rgb = _TEXT
    p.alignment = PP_ALIGN.CENTER
    for para in tf.paragraphs[1:]:
        para.text = ""
    return sh


def _arrow(
    slide,
    x1: Emu,
    y1: Emu,
    x2: Emu,
    y2: Emu,
    *,
    dashed: bool = False,
    label: str | None = None,
) -> None:
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, int(x1), int(y1), int(x2), int(y2)
    )
    conn.line.color.rgb = _EDGE
    conn.line.width = Pt(1.25)
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if label:
        mx = (int(x1) + int(x2)) // 2
        my = (int(y1) + int(y2)) // 2
        tb = slide.shapes.add_textbox(
            Emu(mx - 457200), Emu(my - 114300), Emu(914400), Emu(228600)
        )
        tb.fill.background()
        tb.line.fill.background()
        p = tb.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.name = "Calibri"
        p.font.color.rgb = _EDGE


def _curve(slide, x1, y1, x2, y2) -> None:
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.CURVE, int(x1), int(y1), int(x2), int(y2)
    )
    conn.line.color.rgb = _EDGE
    conn.line.width = Pt(1.75)


def build(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = int(13.3333 * 914400)
    prs.slide_height = int(7.5 * 914400)
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.22), Inches(12.6), Inches(0.55)
    )
    title_box.fill.background()
    title_box.line.fill.background()
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "QAG — per-slot flow + hybrid grading (editable)"
    tp.font.size = Pt(22)
    tp.font.bold = True
    tp.font.name = "Calibri"
    tp.font.color.rgb = _TEXT

    y_main = Inches(1.68)
    y_llm = Inches(2.82)
    w, h = Inches(1.02), Inches(0.78)
    gap = Inches(0.15)

    def cx(i: int):
        return Inches(0.32) + i * (w + gap)

    boxes: list = []
    labels = [
        "Source\ndocument",
        "Generate initial\nquestions (N)",
        "Question for\nthis slot",
        "Answer step\n(≤3 attempts)",
    ]
    for i, lab in enumerate(labels):
        boxes.append(_box(slide, cx(i), y_main, w, h, lab))

    sem = _box(
        slide,
        cx(4),
        y_main,
        w,
        h,
        "Keyword\n(optional)\nor LLM judge",
    )
    llm = _box(
        slide,
        cx(4),
        y_llm,
        w,
        h,
        "LLM judge\n(Qwen, :7101)\nfallback",
    )
    verdict = _box(
        slide,
        cx(5),
        y_main,
        w,
        h,
        "Verdict\n(ground + conf)",
    )
    ok = _box(
        slide,
        cx(6),
        y_main,
        Inches(0.98),
        h,
        "Pass\n(keep pair)",
    )
    save = _box(
        slide,
        cx(7),
        y_main,
        Inches(1.08),
        h,
        "Save\nN pairs +\nsummary",
    )

    rep_y = Inches(3.42)
    rep = _box(
        slide,
        cx(5) + Inches(0.02),
        rep_y,
        Inches(1.28),
        Inches(0.88),
        "Fail:\nregenerate 1 question\n(same slot)",
    )

    for i in range(len(boxes) - 1):
        a = boxes[i]
        b = boxes[i + 1]
        _arrow(
            slide,
            Emu(a.left + a.width),
            Emu(a.top + a.height // 2),
            Emu(b.left),
            Emu(b.top + b.height // 2),
        )

    ans = boxes[3]
    _arrow(
        slide,
        Emu(ans.left + ans.width),
        Emu(ans.top + ans.height // 2),
        Emu(sem.left),
        Emu(sem.top + sem.height // 2),
    )

    _arrow(
        slide,
        Emu(sem.left + sem.width),
        Emu(sem.top + sem.height // 2),
        Emu(verdict.left),
        Emu(verdict.top + verdict.height // 2),
        dashed=True,
        label="clear / high conf",
    )
    _arrow(
        slide,
        Emu(sem.left + sem.width // 2),
        Emu(sem.top + sem.height),
        Emu(llm.left + llm.width // 2),
        Emu(llm.top),
        label="low conf / ungrounded",
    )
    _arrow(
        slide,
        Emu(llm.left + llm.width),
        Emu(llm.top + llm.height // 2),
        Emu(verdict.left),
        Emu(verdict.top + verdict.height // 2),
    )

    _arrow(
        slide,
        Emu(verdict.left + verdict.width),
        Emu(verdict.top + verdict.height // 2),
        Emu(ok.left),
        Emu(ok.top + ok.height // 2),
        label="pass",
    )
    _arrow(
        slide,
        Emu(ok.left + ok.width),
        Emu(ok.top + ok.height // 2),
        Emu(save.left),
        Emu(save.top + save.height // 2),
        dashed=True,
        label="all slots",
    )

    _arrow(
        slide,
        Emu(verdict.left + verdict.width // 2),
        Emu(verdict.top + verdict.height),
        Emu(rep.left + rep.width // 2),
        Emu(rep.top),
        label="fail",
    )

    _curve(
        slide,
        Emu(rep.left),
        Emu(rep.top + rep.height // 2),
        Emu(ans.left + ans.width // 2),
        Emu(ans.top + ans.height),
    )
    tb = slide.shapes.add_textbox(
        Inches(2.75), Inches(2.95), Inches(1.85), Inches(0.35)
    )
    tb.fill.background()
    tb.line.fill.background()
    p = tb.text_frame.paragraphs[0]
    p.text = "new Q → answer again"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = _EDGE

    foot = slide.shapes.add_textbox(
        Inches(0.32), Inches(6.5), Inches(12.65), Inches(0.78)
    )
    foot.fill.background()
    foot.line.fill.background()
    fp = foot.text_frame.paragraphs[0]
    fp.text = (
        "Source: run_qa_pipeline.py; grading via "
        "check_hallucination(method='hybrid') — LLM judge only; "
        "then max_answer_attempts and "
        "generate_questions(num_questions=1) on fail."
    )
    fp.font.size = Pt(9)
    fp.font.name = "Calibri"
    fp.font.color.rgb = _EDGE

    prs.save(str(output))
    print(f"Wrote {output}")


def main() -> None:
    ap = argparse.ArgumentParser(
        description="Build editable slot-flow PPTX."
    )
    default_out = Path(
        "docs/architecture/diagrams/"
        "qag_slide1_slot_flow_editable.pptx"
    )
    ap.add_argument(
        "-o",
        "--output",
        type=Path,
        default=default_out,
        help="Output .pptx path",
    )
    args = ap.parse_args()
    build(args.output.resolve())


if __name__ == "__main__":
    main()
