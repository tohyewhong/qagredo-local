#!/usr/bin/env python3
"""Build a polished deck explaining QAG per-document JSON output fields.

Reads a pipeline analysis JSON (or the bundled sample). Run after a real
pipeline job by passing --analysis path/to/doc_*_analysis.json.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, List, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
DEFAULT_SAMPLE = ROOT / "examples" / "sample_qag_doc_analysis.json"
DEFAULT_OUT = ROOT / "docs" / "QAG_Output_Fields_Overview.pptx"
# Slide 2 prefers the LangGraph node/routing diagram; falls back to the
# high-level pipeline strip if the file is missing.
LANGGRAPH_DETAIL_IMG = (
    ROOT / "docs" / "qag_langgraph_pipeline_detailed.png"
)
FULL_PIPELINE_IMG = ROOT / "docs" / "qag_full_pipeline_flow_16x9.png"

NAVY = RGBColor(15, 44, 96)
NAVY_BRIGHT = RGBColor(33, 83, 160)
SLATE = RGBColor(35, 48, 69)
MUTED = RGBColor(90, 104, 128)
PANEL = RGBColor(246, 250, 255)
CARD_A = RGBColor(232, 242, 255)
CARD_B = RGBColor(255, 252, 245)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(200, 214, 235)


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _full_bleed_bg(slide, rgb: RGBColor) -> None:
    r = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    r.fill.solid()
    r.fill.fore_color.rgb = rgb
    r.line.fill.background()


def _left_accent_bar(slide, width_in: float = 0.22) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(width_in),
        Inches(7.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY_BRIGHT
    bar.line.fill.background()


def _top_rule(slide) -> None:
    rule = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(1.28),
        Inches(12.25),
        Inches(0.04),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = LINE
    rule.line.fill.background()


def _footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.35)
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def _slide_heading(slide, title: str, subtitle: str = "") -> None:
    tbox = slide.shapes.add_textbox(
        Inches(0.55), Inches(0.38), Inches(12.0), Inches(0.72)
    )
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        sbox = slide.shapes.add_textbox(
            Inches(0.55), Inches(0.98), Inches(11.8), Inches(0.42)
        )
        stf = sbox.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = MUTED
    _top_rule(slide)


def _rounded_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = LINE


def _card(
    slide,
    title: str,
    body: str,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = LINE
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = NAVY
    bp = tf.add_paragraph()
    bp.text = body
    bp.font.size = Pt(11)
    bp.font.color.rgb = SLATE
    bp.space_before = Pt(4)
    bp.line_spacing = 1.15


def _bullets_in_slide(
    slide,
    lines: Sequence[str],
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 15,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, text in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"•  {text}"
        para.font.size = Pt(size)
        para.font.color.rgb = SLATE
        para.space_after = Pt(7)
        para.line_spacing = 1.2


def _add_hero(prs: Presentation, doc_id: str, grade: str) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _left_accent_bar(slide)

    kicker = slide.shapes.add_textbox(
        Inches(0.65), Inches(1.35), Inches(11.5), Inches(0.4)
    )
    kf = kicker.text_frame
    kf.clear()
    kp = kf.paragraphs[0]
    kp.text = "QAG · Per-document deliverable"
    kp.font.size = Pt(14)
    kp.font.color.rgb = NAVY_BRIGHT
    kp.font.bold = True

    hbox = slide.shapes.add_textbox(
        Inches(0.65), Inches(1.85), Inches(11.5), Inches(1.5)
    )
    hf = hbox.text_frame
    hf.clear()
    hp = hf.paragraphs[0]
    hp.text = "Output JSON:\nwhat each field means"
    hp.font.size = Pt(40)
    hp.font.bold = True
    hp.font.color.rgb = NAVY
    hp.line_spacing = 1.05

    # Badge pills
    def pill(px: float, py: float, label: str) -> None:
        sh = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(px),
            Inches(py),
            Inches(3.2),
            Inches(0.42),
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = PANEL
        sh.line.color.rgb = LINE
        tf = sh.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = SLATE
        p.alignment = PP_ALIGN.CENTER

    pill(0.65, 3.55, f"Document id · {doc_id}")
    pill(4.1, 3.55, f"Sample grade · {grade}")

    sub = slide.shapes.add_textbox(
        Inches(0.65), Inches(4.35), Inches(11.8), Inches(1.1)
    )
    sf = sub.text_frame
    sf.clear()
    sp = sf.paragraphs[0]
    sp.text = (
        "One analysis file per source document: questions, grounded "
        "answers, citations, hybrid grading, and run timings — ready for "
        "audit and downstream systems."
    )
    sp.font.size = Pt(16)
    sp.font.color.rgb = SLATE
    sp.line_spacing = 1.25

    _footer(slide, "QAG output schema overview")


def _resolve_pipeline_slide_image(
    explicit: Path | None,
) -> tuple[Path | None, str]:
    """Pick diagram file for slide 2; return (path or None, footer note)."""
    chain: List[Path] = []
    if explicit is not None:
        chain.append(explicit)
    chain.extend([LANGGRAPH_DETAIL_IMG, FULL_PIPELINE_IMG])
    for candidate in chain:
        if candidate.is_file():
            if candidate.resolve() == LANGGRAPH_DETAIL_IMG.resolve():
                note = (
                    "Detail flow: "
                    "docs/qag_langgraph_pipeline_detailed.png "
                    "(source: qag_langgraph_pipeline_detailed.dot)"
                )
            else:
                note = (
                    "Overview: "
                    "docs/qag_full_pipeline_flow_16x9.png "
                    "(source: qag_full_pipeline_flow.dot)"
                )
            return candidate, note
    return None, "Add langgraph or full-pipeline PNG under docs/"


def _add_pipeline_slide(
    prs: Presentation,
    img_path: Path | None = None,
) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _slide_heading(
        slide,
        "End-to-end context",
        (
            "LangGraph runtime: each node updates state; END serializes one "
            "doc_*_analysis.json. Fields in this deck map to that state."
        ),
    )
    resolved, footer_src = _resolve_pipeline_slide_image(img_path)
    # Leave room above the footer (7.05) for caption panel.
    pic_top = 1.36
    pic_h = 4.72
    if resolved is not None:
        frame = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.38),
            Inches(pic_top),
            Inches(12.55),
            Inches(pic_h + 0.1),
        )
        frame.fill.solid()
        frame.fill.fore_color.rgb = WHITE
        frame.line.color.rgb = LINE
        slide.shapes.add_picture(
            str(resolved),
            Inches(0.48),
            Inches(pic_top + 0.08),
            Inches(12.35),
            Inches(pic_h),
        )
        cap_y = pic_top + pic_h + 0.14
        _rounded_panel(slide, 0.52, cap_y, 12.28, 0.62, PANEL)
        cap = slide.shapes.add_textbox(
            Inches(0.68),
            Inches(cap_y + 0.08),
            Inches(11.95),
            Inches(0.46),
        )
        ctf = cap.text_frame
        ctf.clear()
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = (
            "Detail flow — question generation (with validation loops), "
            "answer generation, hybrid grade_primary, optional "
            "grade_fallback_llm when routing rules fire; then writer emits "
            "document, qa_pairs (incl. hallucination_check per slot), "
            "generation blocks, grading_summary, run_metrics."
        )
        cp.font.size = Pt(10)
        cp.font.color.rgb = SLATE
        cp.line_spacing = 1.12
    else:
        _bullets_in_slide(
            slide,
            [
                "Missing diagram assets. Add either:",
                "docs/qag_langgraph_pipeline_detailed.png, or",
                "docs/qag_full_pipeline_flow_16x9.png",
            ],
            0.65,
            1.55,
            11.5,
            2.4,
            14,
        )
    _footer(slide, footer_src)


def _add_keys_grid(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _slide_heading(
        slide,
        "Six top-level keys",
        "One timestamped folder per run; one JSON file per document.",
    )
    items: List[Tuple[str, str, RGBColor]] = [
        (
            "document",
            "Normalized snapshot: id, title, source, type, content.",
            CARD_A,
        ),
        (
            "qa_pairs",
            "Each slot: question, answer, hallucination_check, "
            "citation_spans, provenance.",
            CARD_B,
        ),
        (
            "question_generation",
            "Model, types, timestamps; optional validation detail.",
            CARD_A,
        ),
        (
            "answer_generation",
            "Answer model, provider, timestamps, count.",
            CARD_B,
        ),
        (
            "grading_summary",
            "Letter grade, confidence, method, judge model.",
            CARD_A,
        ),
        (
            "run_metrics",
            "Stage timings (s) and quality counters.",
            CARD_B,
        ),
    ]
    # 3 columns × 2 rows
    cw, ch, gx, gy = 3.95, 1.38, 0.55, 1.42
    col_gap, row_gap = 0.35, 0.18
    positions = [
        (gx, gy),
        (gx + cw + col_gap, gy),
        (gx + 2 * (cw + col_gap), gy),
        (gx, gy + ch + row_gap),
        (gx + cw + col_gap, gy + ch + row_gap),
        (gx + 2 * (cw + col_gap), gy + ch + row_gap),
    ]
    for i, (title, body, fill) in enumerate(items):
        x, y = positions[i]
        _card(slide, title, body, x, y, cw, ch, fill)
    _footer(slide, "Matches run_qa_pipeline combined_result")


def _add_document_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _slide_heading(
        slide,
        "document",
        "Single source of truth for grading text.",
    )
    _rounded_panel(slide, 0.52, 1.38, 6.15, 5.35, PANEL)
    _bullets_in_slide(
        slide,
        [
            "id — from input or derived",
            "title — when provided",
            "source / type — optional provenance",
            "content — plain text (aliases merged: "
            "content, text, body, …)",
        ],
        0.72,
        1.58,
        5.75,
        4.9,
        15,
    )
    _rounded_panel(slide, 6.85, 1.38, 5.95, 5.35, CARD_A)
    note = slide.shapes.add_textbox(
        Inches(7.05), Inches(1.58), Inches(5.55), Inches(4.9)
    )
    nf = note.text_frame
    nf.clear()
    nf.word_wrap = True
    p = nf.paragraphs[0]
    p.text = "Why it matters"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = NAVY
    for line in (
        "Auditors and tools always see the exact text the pipeline used — "
        "even when the JSONL stored prose under another key.",
        "",
        "Hallucination checks and citation offsets refer to this content.",
    ):
        bp = nf.add_paragraph()
        bp.text = line
        bp.font.size = Pt(13)
        bp.font.color.rgb = SLATE
        bp.line_spacing = 1.22
    _footer(slide, "document snapshot")


def _add_qa_slide(prs: Presentation, data: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _slide_heading(
        slide,
        "qa_pairs[]",
        "One object per question slot after retries complete.",
    )
    _rounded_panel(slide, 0.52, 1.35, 6.0, 5.42, PANEL)
    _bullets_in_slide(
        slide,
        [
            "question, answer — final strings after slot retries",
            "hallucination_check — right after answer: full verifier payload "
            "(is_grounded, confidence, method, issues, llm_verdict, "
            "grounding_why when enabled)",
            "citation_spans — list of {start, end, text} char offsets into "
            "document.content (from supporting_evidence, not stored in JSON)",
            "citation_notes — fragments not found verbatim after the same "
            "split / list-prefix strip / dedupe as spans",
            "source_doc_id, source_title — trace to input row / title",
        ],
        0.72,
        1.52,
        5.65,
        5.1,
        14,
    )
    _rounded_panel(slide, 6.68, 1.35, 6.12, 5.42, CARD_B)
    right_title = slide.shapes.add_textbox(
        Inches(6.88), Inches(1.48), Inches(5.75), Inches(0.38)
    )
    rt = right_title.text_frame
    rt.clear()
    rp = rt.paragraphs[0]
    rp.text = "Live excerpt (first pair)"
    rp.font.bold = True
    rp.font.size = Pt(13)
    rp.font.color.rgb = NAVY
    pairs = data.get("qa_pairs")
    if isinstance(pairs, list) and pairs and isinstance(pairs[0], dict):
        first = pairs[0]
        q = str(first.get("question", ""))[:220]
        a = str(first.get("answer", ""))[:220]
        g = first.get("hallucination_check")
        if not isinstance(g, dict):
            g = first.get("grading")
        meta = ""
        if isinstance(g, dict):
            meta = (
                f"grounded={g.get('is_grounded')} · "
                f"conf={g.get('confidence')} · "
                f"{g.get('method', '')}"
            )
        body_lines = [f"Q: {q}", "", f"A: {a}", "", meta]
    else:
        body_lines = ["(No qa_pairs in sample file.)"]
    body_box = slide.shapes.add_textbox(
        Inches(6.88), Inches(1.95), Inches(5.75), Inches(4.65)
    )
    bf = body_box.text_frame
    bf.clear()
    bf.word_wrap = True
    for i, line in enumerate(body_lines):
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        para.text = line
        para.font.size = Pt(11)
        para.font.color.rgb = SLATE
        para.line_spacing = 1.15
    _footer(slide, "Per-slot QA record")


def _add_grading_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _slide_heading(
        slide,
        "hallucination_check (per slot)",
        "Same fields grade_qa_results puts in check_result.",
    )
    cols = 2
    headers = ("Field", "Role")
    data_rows = [
        ("is_grounded", "Pass/fail vs. confidence + sentence rules."),
        ("confidence", "0.0–1.0 aggregate for this answer."),
        ("method", "semantic, keyword, llm, hybrid variants, fallbacks."),
        ("issues", "Human-readable problems or judge notes."),
        (
            "grounded_sentences / ungrounded_sentences",
            "Sentence-level split for debugging (use len() if you need "
            "counts).",
        ),
        ("llm_verdict", "When judge runs: verdict, confidence, reason."),
        (
            "grounding_why",
            "If enabled: brief why when grounded but citation_spans/notes "
            "empty (llm_verdict.reason or extra judge call).",
        ),
    ]
    rows = 1 + len(data_rows)
    tw = Inches(12.0)
    th = Inches(5.35)
    left = Inches(0.65)
    top = Inches(1.38)
    tbl = slide.shapes.add_table(rows, cols, left, top, tw, th).table
    tbl.columns[0].width = Inches(3.15)
    tbl.columns[1].width = Inches(8.85)
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_BRIGHT
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
    for r, (fld, role) in enumerate(data_rows, start=1):
        c0 = tbl.cell(r, 0)
        c1 = tbl.cell(r, 1)
        c0.text = fld
        c1.text = role
        fill = PANEL if r % 2 == 1 else WHITE
        for cell in (c0, c1):
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = SLATE
    _footer(slide, "utils/hallucination_checker.py")


def _col_heading(slide, x: float, y: float, w: float, text: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(0.35)
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY


def _add_grounding_vs_citations_slide(prs: Presentation) -> None:
    """Explain is_grounded vs citation_spans — independent audit signals."""
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _slide_heading(
        slide,
        "Two different checks (do not conflate them)",
        (
            "is_grounded grades the answer against the full document. "
            "citation_* locates quoted evidence strings in the document."
        ),
    )
    _col_heading(
        slide,
        0.58,
        1.32,
        5.85,
        "A. qa_pairs[].hallucination_check (is_grounded)",
    )
    _bullets_in_slide(
        slide,
        [
            "What is scored: the answer text (and optionally the "
            "question as judge context), compared to document.content.",
            "Producer: grade_qa_results → check_hallucination in "
            "utils/hallucination_checker.py.",
            "Semantic/keyword: each answer sentence vs document chunks or "
            "phrases; ratio + rules → confidence; is_grounded requires "
            "no ungrounded sentences and confidence ≥ 0.7 (fixed in code).",
            "LLM/hybrid: judge reads truncated document + Q + A; verdict "
            "SUPPORTED plus judge confidence ≥ 0.7.",
            "Purpose: “Did the model invent facts relative to the doc?”",
        ],
        0.58,
        1.72,
        5.85,
        5.25,
        12,
    )
    _col_heading(slide, 6.72, 1.32, 6.0, "B. citation_spans / citation_notes")
    _bullets_in_slide(
        slide,
        [
            "What is scored: supporting_evidence for that slot only "
            "(parallel list from answer generation — not written into JSON).",
            "Producer: run_qa_pipeline.build_qa_pairs → "
            "_evidence_to_citation_spans on document.content.",
            "Algorithm: split (newlines, semicolons); strip - * • and 1. / 1) "
            "prefixes; dedupe identical fragments; search verbatim (or "
            "whitespace-relaxed); hit → {start,end,text}; miss → "
            "citation_notes.",
            "Purpose: “Which exact spans did the model claim it quoted?”",
            "Empty spans + empty notes usually means empty or unusable "
            "supporting_evidence — not “ungrounded” by the grader.",
        ],
        6.72,
        1.72,
        6.05,
        5.25,
        12,
    )
    _rounded_panel(slide, 0.52, 6.12, 12.28, 0.95, PANEL)
    _bullets_in_slide(
        slide,
        [
            "You can explain is_grounded without citations: verifier never "
            "reads supporting_evidence.",
            "You can have is_grounded true and citation_spans [] if the "
            "answer paraphrases well but evidence quotes are missing or "
            "non-verbatim.",
            "Teach both: verifier = holistic support; citations = explicit "
            "quote alignment for auditors.",
        ],
        0.7,
        6.22,
        11.95,
        0.82,
        11,
    )
    _footer(slide, "run_qa_pipeline.py + hallucination_checker.py")


def _add_citation_pipeline_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _slide_heading(
        slide,
        "citation_spans & citation_notes — step by step",
        "Same character offsets as document.content in the output file.",
    )
    _bullets_in_slide(
        slide,
        [
            "Index i: qa_pairs[i] aligns with supporting_evidence[i] inside "
            "the internal qa_result (answers and evidence lists same length).",
            "Split: newlines, then semicolons; strip leading - * • and "
            "numeric list prefixes; dedupe by collapsed whitespace; drop "
            "fragments under 3 chars (_split_evidence_fragments).",
            "Match: first exact substring of fragment in content; else "
            "regex with flexible whitespace between tokens "
            "(_find_quote_span).",
            "Hit: append {start, end, text} where text = content[start:end].",
            "Miss: append fragment string to citation_notes (paraphrase, "
            "hallucinated quote, or typo).",
            "If supporting_evidence[i] is blank, both arrays are empty — "
            "grading may still pass.",
        ],
        0.58,
        1.38,
        11.85,
        5.55,
        13,
    )
    _rounded_panel(slide, 0.52, 6.05, 12.28, 1.02, CARD_A)
    _bullets_in_slide(
        slide,
        [
            "Console: [WARN] No citation_spans resolved when evidence "
            "non-empty but nothing matched (paraphrase).",
            "Downstream UI: highlight document.content[start:end] for each "
            "span; show citation_notes as “could not locate in doc”.",
        ],
        0.68,
        6.18,
        11.95,
        0.88,
        11,
    )
    _footer(slide, "_evidence_to_citation_spans")


def _add_is_grounded_methods_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _slide_heading(
        slide,
        "is_grounded & confidence — by hallucination.method",
        "config.<profile>.yaml → hallucination.method; answer gate uses "
        "answer_generation.min_confidence_threshold (default 0.7).",
    )
    nmeth = 4
    rows = 1 + nmeth
    cols = 2
    tbl = slide.shapes.add_table(
        rows, cols, Inches(0.58), Inches(1.38), Inches(12.15), Inches(4.85)
    ).table
    tbl.columns[0].width = Inches(2.35)
    tbl.columns[1].width = Inches(9.8)
    hdr = ("Method", "How confidence & is_grounded are set")
    data_rows = [
        (
            "keyword",
            "Per-sentence: key phrases in doc (or generic/hedge heuristics). "
            "confidence = grounded/total sentences; +0.2 cap 1.0 if answer "
            "contains e.g. “not enough information”. is_grounded = confidence "
            "≥ 0.7 AND ungrounded_sentences empty.",
        ),
        (
            "semantic",
            "Keyword / LLM judge: overlap or judge verdict; "
            "max cosine ≥ 0.5 ⇒ grounded. confidence = grounded/total. "
            "Same is_grounded rule as keyword (≥ 0.7, no ungrounded).",
        ),
        (
            "llm",
            "Single judge call; JSON verdict SUPPORTED / NOT_SUPPORTED + "
            "confidence. is_grounded = SUPPORTED AND confidence ≥ 0.7. "
            "Full answer + truncated doc in prompt.",
        ),
        (
            "hybrid",
            "Semantic first; if any ungrounded sentences, LLM re-judges "
            "whole answer. Override: is_grounded true if LLM supported & "
            "conf ≥ 0.7. Else min(semantic, llm) confidence and semantic "
            "ungrounded kept.",
        ),
    ]
    for c, h in enumerate(hdr):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_BRIGHT
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
    for r, (meth, expl) in enumerate(data_rows, start=1):
        c0 = tbl.cell(r, 0)
        c1 = tbl.cell(r, 1)
        c0.text = meth
        c1.text = expl
        fill = PANEL if r % 2 == 1 else WHITE
        for cell in (c0, c1):
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(10)
                para.font.color.rgb = SLATE
    _footer(slide, "Threshold 0.7 is policy in code, not computed from doc")


def _add_question_validation_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _slide_heading(
        slide,
        "question_generation — fields & hidden validation",
        "Same check_hallucination API, but the “answer” argument is the "
        "question string.",
    )
    _bullets_in_slide(
        slide,
        [
            "model, provider, timestamp, timezone — who/when generated "
            "questions.",
            "num_questions, complexity, question_types — prompt shaping "
            "and presets.",
            "question_validation — list per slot when enable_rejection or "
            "enable_comprehensiveness_check: original_question, "
            "final_question, validation_info {is_grounded, confidence, "
            "issues, attempts, was_regenerated}, optional "
            "comprehensiveness_check.",
            "validation uses question_generation.validation.method "
            "(default semantic); checks question text vs document — NOT "
            "the later answer grading.",
        ],
        0.58,
        1.38,
        11.85,
        5.35,
        13,
    )
    _rounded_panel(slide, 0.52, 6.05, 12.28, 1.02, PANEL)
    _bullets_in_slide(
        slide,
        [
            "Story for stakeholders: “We optionally reject/regenerate "
            "questions that are not about the document before any answer "
            "is written.”",
        ],
        0.68,
        6.2,
        11.9,
        0.82,
        12,
    )
    _footer(slide, "utils/question_generator.py")


def _add_answer_generation_gate_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _slide_heading(
        slide,
        "answer_generation — fields & regeneration gate",
        "During generation, answers can be rewritten until grounding passes.",
    )
    _bullets_in_slide(
        slide,
        [
            "model, provider, timestamp, timezone, num_answers — batch "
            "answer pass metadata.",
            "Per-answer validation (optional): check_hallucination with "
            "config hallucination.method; min_confidence_threshold from "
            "answer_generation in config.",
            "Final JSON verifier (qa_pairs.hallucination_check) comes "
            "from a separate grade_qa_results pass on the finished "
            "qa_result — usually matches generation-time check if method "
            "unchanged.",
            "supporting_evidence is produced with the answer; used only to "
            "build citation_*; omit from export by design.",
        ],
        0.58,
        1.38,
        11.85,
        5.55,
        13,
    )
    _footer(slide, "utils/answer_generator.py + run_qa_pipeline slot loop")


def _add_hallucination_checks_detail_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _slide_heading(
        slide,
        "No top-level hallucination_checks",
        "Verifier payload per slot: qa_pairs[].hallucination_check.",
    )
    _bullets_in_slide(
        slide,
        [
            "combined_result no longer emits a duplicate list; use "
            "qa_pairs[i].hallucination_check only.",
            "grading_summary still aggregates from those per-slot objects.",
            "Legacy files may still have key grading or a top-level "
            "hallucination_checks array — readers should accept both.",
        ],
        0.58,
        1.38,
        11.85,
        5.55,
        13,
    )
    _footer(slide, "run_qa_pipeline combined_result")


def _add_rollups_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _slide_heading(slide, "Run-level rollups", "One-line health + timing.")
    _card(
        slide,
        "grading_summary",
        "overall_grade (A–F), overall_confidence, grading_method "
        "(e.g. hybrid or average_of_each_qa_pair), judge_model.",
        0.55,
        1.42,
        5.95,
        2.45,
        CARD_A,
    )
    _card(
        slide,
        "run_metrics",
        "timings_seconds: question_generation, answer_generation, "
        "grading. quality_counters: retries / rewrites when enabled.",
        6.68,
        1.42,
        5.95,
        2.45,
        CARD_B,
    )
    _rounded_panel(slide, 0.55, 4.05, 12.08, 2.75, PANEL)
    _bullets_in_slide(
        slide,
        [
            "Per-slot hallucination_check is next to answer in qa_pairs.",
            "Use grading_summary for dashboards; drill into qa_pairs[] "
            "for evidence.",
        ],
        0.75,
        4.28,
        11.7,
        2.35,
        14,
    )
    _footer(slide, "Aggregate vs. per-slot detail")


def _add_closing(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _left_accent_bar(slide)
    _slide_heading(slide, "Files & regeneration", "")
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.55),
        Inches(1.45),
        Inches(12.15),
        Inches(4.85),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = LINE
    cmds = [
        "Output path pattern:",
        "output/<provider>/<model>/<timestamp>/doc_<id>_analysis.json",
        "",
        "This deck adds deep slides: grounding vs citations, citation "
        "algorithm, is_grounded by method, question/answer metadata, "
        "per-slot hallucination_check.",
        "",
        "Rebuild this deck:",
        "python3 scripts/utils/build_qag_output_fields_ppt.py \\",
        "  --analysis path/to/doc_*_analysis.json",
        "",
        "Run summary (all docs in folder):",
        "bash run.sh --summarize --latest --json",
    ]
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(12)
    tf.margin_top = tf.margin_bottom = Pt(10)
    for i, line in enumerate(cmds):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(13 if i in (0, 3, 7) else 12)
        if i in (0, 3, 7):
            para.font.bold = True
            para.font.color.rgb = NAVY
        else:
            para.font.color.rgb = SLATE
        if line.startswith(("python3", " ", "output", "bash")):
            para.font.name = "Consolas"
        para.space_after = Pt(4)
    _footer(slide, "QAG")


def _load_analysis(path: Path) -> Dict[str, Any]:
    with open(path, encoding="utf-8") as f:
        data = json.load(f)
    if not isinstance(data, dict):
        raise ValueError("Analysis file must be a JSON object")
    return data


def build_deck(
    analysis: Dict[str, Any],
    out_path: Path,
    *,
    include_pipeline_image: bool = True,
    pipeline_slide_image: Path | None = None,
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    doc = analysis.get("document", {})
    doc_id = doc.get("id", "—") if isinstance(doc, dict) else "—"
    gs = analysis.get("grading_summary", {})
    og = gs.get("overall_grade", "—") if isinstance(gs, dict) else "—"

    _add_hero(prs, str(doc_id), str(og))
    if include_pipeline_image:
        _add_pipeline_slide(prs, pipeline_slide_image)
    _add_keys_grid(prs)
    _add_document_slide(prs)
    _add_qa_slide(prs, analysis)
    _add_grading_slide(prs)
    _add_grounding_vs_citations_slide(prs)
    _add_citation_pipeline_slide(prs)
    _add_is_grounded_methods_slide(prs)
    _add_question_validation_slide(prs)
    _add_answer_generation_gate_slide(prs)
    _add_hallucination_checks_detail_slide(prs)
    _add_rollups_slide(prs)
    _add_closing(prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Build QAG output-field overview slides."
    )
    parser.add_argument(
        "--analysis",
        type=Path,
        default=DEFAULT_SAMPLE,
        help="Per-document analysis JSON (default: examples sample)",
    )
    parser.add_argument(
        "--out",
        type=Path,
        default=DEFAULT_OUT,
        help="Output .pptx path",
    )
    parser.add_argument(
        "--no-pipeline-image",
        action="store_true",
        help="Omit pipeline diagram slide if present",
    )
    parser.add_argument(
        "--pipeline-slide-image",
        type=Path,
        default=None,
        help="Override slide-2 diagram (default: langgraph detail, then "
        "full-pipeline PNG)",
    )
    args = parser.parse_args()
    data = _load_analysis(args.analysis)
    build_deck(
        data,
        args.out,
        include_pipeline_image=not args.no_pipeline_image,
        pipeline_slide_image=args.pipeline_slide_image,
    )
    print(f"Wrote {args.out}")


if __name__ == "__main__":
    main()
