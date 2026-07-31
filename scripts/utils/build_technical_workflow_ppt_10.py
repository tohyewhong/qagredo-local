from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


DOCS_DIR = Path(__file__).resolve().parents[2] / "docs"
ASSET_DIR = DOCS_DIR / "assets"
OUTFILE = DOCS_DIR / "QAG_Technical_Workflow_10slides.pptx"

IMG = {
    "hero": ASSET_DIR / "slide_visual_01_title_hero.png",
    "problem": ASSET_DIR / "slide_visual_02_problem_context.png",
    "pipeline": DOCS_DIR / "qag_full_pipeline_flow_16x9.png",
    "arch": ASSET_DIR / "slide_visual_04_architecture_containers.png",
    "qgen": ASSET_DIR / "slide_visual_05_question_generation.png",
    "aval": ASSET_DIR / "slide_visual_06_answer_validation.png",
    "hybrid": ASSET_DIR / "slide_visual_07_hybrid_grading.png",
    "audit": ASSET_DIR / "slide_visual_08_output_audit.png",
    "tradeoff": ASSET_DIR / "slide_visual_10_tradeoff_balance.png",
    "nrrt": ASSET_DIR / "slide_visual_15_need_risk_response_target.png",
}


def add_title_box(slide, title, subtitle=""):
    t = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.9)
    ).text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(18, 47, 95)
    if subtitle:
        st = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.5)
        ).text_frame
        st.clear()
        sp = st.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(15)
        sp.font.color.rgb = RGBColor(66, 82, 108)


def add_bullets(slide, items, x=0.7, y=1.7, w=6.2, h=4.8, font_size=16):
    tf = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    ).text_frame
    tf.word_wrap = True
    tf.clear()
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(35, 48, 69)
        p.space_after = Pt(8)


def add_image(slide, path, x, y, w, h):
    if path.exists():
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.color.rgb = RGBColor(220, 228, 240)
        slide.shapes.add_picture(
            str(path), Inches(x), Inches(y), Inches(w), Inches(h)
        )


def add_stage_box(
    slide,
    x,
    y,
    w,
    h,
    title,
    lines,
    fill_rgb,
    title_font_size=14,
    body_font_size=11,
):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*fill_rgb)
    box.line.color.rgb = RGBColor(160, 175, 205)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(title_font_size)
    p.font.color.rgb = RGBColor(25, 45, 85)
    for line in lines:
        bp = tf.add_paragraph()
        bp.text = line
        bp.font.size = Pt(body_font_size)
        bp.font.color.rgb = RGBColor(45, 58, 80)


def add_overlay_label(
    slide,
    x,
    y,
    w,
    h,
    text,
    fill_rgb=(255, 255, 255),
    font_rgb=(25, 45, 85),
    font_size=11,
):
    label = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    label.fill.solid()
    label.fill.fore_color.rgb = RGBColor(*fill_rgb)
    label.fill.transparency = 8
    label.line.color.rgb = RGBColor(140, 160, 195)
    tf = label.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*font_rgb)
    p.alignment = PP_ALIGN.CENTER


def add_takeaway(slide, text):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(6.75),
        Inches(12.0),
        Inches(0.45),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(240, 246, 255)
    box.line.color.rgb = RGBColor(170, 190, 225)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Key takeaway: {text}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 56, 106)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1) Title
    s = prs.slides.add_slide(blank)
    add_image(s, IMG["hero"], 0, 0, 13.333, 7.5)
    overlay = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
    overlay.fill.transparency = 18
    overlay.line.fill.background()
    tf = s.shapes.add_textbox(
        Inches(0.8), Inches(1.6), Inches(11.8), Inches(2.2)
    ).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "QAG Technical Workflow + Framework Upgrade"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(18, 47, 95)
    p = tf.add_paragraph()
    p.text = "Balanced view: workflow, algorithm, and tradeoffs"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(66, 82, 108)

    # 2) Why
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "2. Why This Workflow Exists",
        "Problem, design response, target outcome",
    )
    add_bullets(
        s,
        [
            "Need: complex Q&A grounded strictly in source documents.",
            "Risk: raw LLM output can be shallow, repetitive, or hallucinatory.",  # noqa: E501
            "Response: staged generation + validation + adaptive grading fallback.",  # noqa: E501
            "Target: explainable, auditable, high-confidence results.",
        ],
        w=6.3,
        font_size=17,
    )
    add_image(s, IMG["nrrt"], 7.0, 1.7, 6.0, 4.9)
    add_overlay_label(
        s,
        7.15,
        5.95,
        1.25,
        0.36,
        "Need",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
        font_size=10,
    )
    add_overlay_label(
        s,
        8.55,
        5.95,
        1.25,
        0.36,
        "Risk",
        fill_rgb=(255, 238, 225),
        font_rgb=(118, 52, 18),
        font_size=10,
    )
    add_overlay_label(
        s,
        9.95,
        5.95,
        1.55,
        0.36,
        "Response",
        fill_rgb=(233, 252, 244),
        font_rgb=(14, 86, 48),
        font_size=10,
    )
    add_overlay_label(
        s,
        11.65,
        5.95,
        1.15,
        0.36,
        "Target",
        fill_rgb=(240, 236, 255),
        font_rgb=(58, 36, 120),
        font_size=10,
    )
    add_takeaway(
        s, "Quality comes from explicit gates, not from one model call."
    )

    # 3) End-to-end workflow
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "3. End-to-End Workflow",
        "Full pipeline from input to audited output",
    )
    add_image(s, IMG["pipeline"], 0.75, 1.35, 11.95, 3.35)
    # Keep the full-flow image clean; explain via compact lane cards below.
    add_stage_box(
        s,
        0.9,
        4.95,
        3.95,
        1.35,
        "Lane A: Input + Prep",
        [
            "Converter CLI: --input-type routing",
            "Normalization + optional semantic enrich",
        ],
        (232, 241, 255),
        title_font_size=16,
        body_font_size=12,
    )
    add_stage_box(
        s,
        4.95,
        4.95,
        3.95,
        1.35,
        "Lane B: Q/A Quality Loop",
        [
            "LangChain prompt/template/parsing path",
            "Grounding, retries, and coverage checks",
        ],
        (233, 252, 244),
        title_font_size=16,
        body_font_size=12,
    )
    add_stage_box(
        s,
        9.0,
        4.95,
        3.25,
        1.35,
        "Lane C: Grading + Output",
        ["Hybrid scoring path", "Grade, issues, run summary + run metrics"],
        (240, 236, 255),
        title_font_size=16,
        body_font_size=12,
    )
    add_takeaway(
        s,
        "Single visual, three lanes: preparation, quality loop, and final audit output.",  # noqa: E501
    )

    # 4) Architecture
    s = prs.slides.add_slide(blank)
    add_title_box(
        s, "4. Runtime Architecture", "Three-container responsibility split"
    )
    add_image(s, IMG["arch"], 0.6, 1.4, 7.2, 5.4)
    add_stage_box(
        s,
        8.1,
        1.8,
        4.5,
        1.3,
        "Generator (qag-vllm)",
        [
            "Llama | GPU 0 | Port ${VLLM_HOST_PORT}",
            "Question + answer generation",
        ],
        (232, 241, 255),
    )
    add_stage_box(
        s,
        8.1,
        3.3,
        4.5,
        1.3,
        "Judge (qag-vllm-judge)",
        [
            "Qwen | GPU 1 | Port ${VLLM_JUDGE_HOST_PORT}",
            "Independent grounding verification",
        ],
        (240, 236, 255),
    )
    add_stage_box(
        s,
        8.1,
        4.8,
        4.5,
        1.3,
        "Runner (qag-runner)",
        ["CPU | internal", "LangGraph orchestration + grading checks"],
        (233, 252, 244),
    )
    add_takeaway(
        s, "Generator and judge split improves reliability and throughput."
    )

    # 5) Question strategy
    s = prs.slides.add_slide(blank)
    add_title_box(
        s, "5. Question Generation Strategy", "10 types + few-shot effect"
    )
    add_image(s, IMG["qgen"], 0.7, 1.5, 6.0, 4.9)
    add_bullets(
        s,
        [
            "10 types cover analysis -> counterfactual reasoning.",
            "Few-shot uses GOOD/BAD examples to avoid trivial lookup questions.",  # noqa: E501
            "Dedup + grounding + comprehensiveness filters weak questions.",
            "Better upstream questions reduce downstream answer retries.",
        ],
        x=7.0,
        y=1.8,
        w=5.6,
        h=4.8,
        font_size=15,
    )
    add_takeaway(
        s, "Question quality is the biggest lever for final answer quality."
    )

    # 6) Answer generation + validation loop
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "6. Answer Generation + Validation",
        "Structured answer and retry control",
    )
    add_image(s, IMG["aval"], 6.6, 1.5, 6.4, 5.2)
    add_overlay_label(
        s,
        6.8,
        1.85,
        2.0,
        0.38,
        "Generated Answer",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
    )
    add_overlay_label(
        s,
        8.95,
        3.15,
        1.7,
        0.38,
        "Grounding Check",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
    )
    add_overlay_label(
        s,
        10.75,
        1.95,
        2.15,
        0.38,
        "Supported Evidence",
        fill_rgb=(225, 255, 238),
        font_rgb=(14, 86, 48),
    )
    add_overlay_label(
        s,
        10.75,
        3.55,
        2.15,
        0.38,
        "Ungrounded Claims",
        fill_rgb=(255, 238, 225),
        font_rgb=(118, 52, 18),
    )
    add_bullets(
        s,
        [
            "Prompt enforces: answer-only-from-document + supporting evidence.",  # noqa: E501
            "If grounding fails, regenerate up to configured retry limit.",
            "Best attempt is retained with reasons for transparency.",
        ],
        x=0.7,
        y=1.9,
        w=5.6,
        h=4.8,
        font_size=15,
    )
    add_takeaway(
        s, "Retries are controlled and evidence-backed, not open-ended."
    )

    # 7) Coverage validator
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "7. Coverage Validator + Rewrite",
        "Fixes grounded-but-incomplete answers",
    )
    add_image(s, IMG["aval"], 6.7, 1.6, 6.0, 5.0)
    add_overlay_label(
        s,
        6.9,
        1.95,
        1.9,
        0.38,
        "Initial Answer",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
    )
    add_overlay_label(
        s,
        8.9,
        3.0,
        1.9,
        0.38,
        "Coverage Check",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
    )
    add_overlay_label(
        s,
        10.9,
        3.55,
        1.75,
        0.38,
        "Low Coverage",
        fill_rgb=(255, 238, 225),
        font_rgb=(118, 52, 18),
    )
    add_overlay_label(
        s,
        8.75,
        5.05,
        2.2,
        0.38,
        "Targeted Rewrite (x1)",
        fill_rgb=(255, 247, 220),
        font_rgb=(96, 73, 15),
    )
    add_bullets(
        s,
        [
            "Detects when answer misses parts of the question.",
            "Runs one targeted rewrite using missing-point feedback.",
            "Rewrite is accepted only if grounding still passes.",
        ],
        x=0.7,
        y=1.9,
        w=5.6,
        h=4.8,
        font_size=15,
    )
    add_takeaway(
        s, "Completeness improves without raising hallucination risk."
    )

    # 8) Hybrid grading
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "8. Hybrid Grading Decision Path",
        "Speed-first with robust fallback",
    )
    add_image(s, IMG["hybrid"], 0.7, 1.5, 6.4, 5.1)
    add_bullets(
        s,
        [
            "Path A: semantic confidence high -> accept quickly.",
            "Path B: low confidence -> invoke independent LLM judge.",
            "Final score combines evidence for confidence and grade.",
            "Hybrid is default for speed/quality balance.",
        ],
        x=7.3,
        y=1.9,
        w=5.4,
        h=4.8,
        font_size=15,
    )
    add_takeaway(s, "Only uncertain cases pay the judge-model latency cost.")

    # 9) Output schema
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "9. Output Schema Made Simple",
        "How to read one result in 30 seconds",
    )
    add_image(s, IMG["audit"], 8.0, 1.5, 4.7, 2.5)
    add_stage_box(
        s,
        0.7,
        1.8,
        3.8,
        2.0,
        "Block 1: Context",
        ["document.id/title/source", "traceability"],
        (233, 242, 255),
    )
    add_stage_box(
        s,
        4.85,
        1.8,
        3.8,
        2.0,
        "Block 2: QA details",
        ["question + answer + grading", "per-pair issues"],
        (233, 252, 244),
    )
    add_stage_box(
        s,
        0.7,
        4.1,
        3.8,
        2.0,
        "Block 3: Generation",
        ["question/answer metadata", "model + timestamps"],
        (240, 236, 255),
    )
    add_stage_box(
        s,
        4.85,
        4.1,
        3.8,
        2.0,
        "Block 4: Summary",
        ["overall grade/confidence", "run metadata + semantic_enrichment"],
        (255, 240, 229),
    )
    add_bullets(
        s,
        [
            "If grade is low, inspect QA issues first, then trace generation metadata."  # noqa: E501
        ],
        x=8.0,
        y=4.25,
        w=4.8,
        h=1.6,
        font_size=12,
    )
    add_takeaway(s, "Outputs are built for fast audit, not just storage.")

    # 10) Tradeoffs + recommendations
    s = prs.slides.add_slide(blank)
    add_title_box(
        s, "10. Tradeoff Summary + Recommendations", "Senior decision view"
    )
    add_image(s, IMG["tradeoff"], 0.8, 1.5, 5.8, 5.0)
    # In-image labels to make the tradeoff visual self-explanatory.
    add_overlay_label(
        s,
        1.15,
        5.95,
        1.65,
        0.36,
        "Accuracy",
        fill_rgb=(225, 255, 238),
        font_rgb=(14, 86, 48),
        font_size=10,
    )
    add_overlay_label(
        s,
        3.0,
        5.95,
        1.65,
        0.36,
        "Latency",
        fill_rgb=(255, 247, 220),
        font_rgb=(96, 73, 15),
        font_size=10,
    )
    add_overlay_label(
        s,
        4.85,
        5.95,
        1.35,
        0.36,
        "Cost",
        fill_rgb=(240, 236, 255),
        font_rgb=(58, 36, 120),
        font_size=10,
    )
    add_overlay_label(
        s,
        2.35,
        2.0,
        2.6,
        0.36,
        "Balance, not optimize one only",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
        font_size=9,
    )
    add_stage_box(
        s,
        6.9,
        1.8,
        5.7,
        3.7,
        "Decision Tradeoffs",
        [
            "Higher depth checks -> better quality, higher latency",
            "Lower temperature -> more factual, less linguistic flexibility",
            "Hybrid grading -> better speed/quality, added routing logic",
            "Coverage rewrite -> better completeness, one extra call",
        ],
        (245, 248, 255),
    )
    add_bullets(
        s,
        [
            "Image mapping: Accuracy (quality checks), Latency (retry/fallback depth), Cost (judge-model usage).",  # noqa: E501
            "Decision: keep current architecture; tune complexity and retry thresholds first.",  # noqa: E501
            "Deployment: use 6-file split-model (separate llama/qwen archives + bundle updates).",  # noqa: E501
            "Operations: add observability metrics and fast/standard/strict presets.",  # noqa: E501
        ],
        x=6.9,
        y=5.7,
        w=5.7,
        h=1.0,
        font_size=12,
    )
    add_takeaway(
        s, "Operate with explicit, measurable speed-quality tradeoffs."
    )

    OUTFILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTFILE))
    print(str(OUTFILE))


if __name__ == "__main__":
    build()
