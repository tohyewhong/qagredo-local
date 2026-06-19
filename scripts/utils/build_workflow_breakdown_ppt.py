from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


DOCS_DIR = Path(__file__).resolve().parents[2] / "docs"
OUTFILE = DOCS_DIR / "QAGRedo_End_to_End_Workflow_Breakdown.pptx"

IMG = {
    "full": DOCS_DIR / "qagredo_full_pipeline_flow_16x9.png",
    "part1": DOCS_DIR / "qagredo_workflow_part1_input_prep_16x9.png",
    "part2": DOCS_DIR / "qagredo_workflow_part2_qa_loop_16x9.png",
    "part3": DOCS_DIR / "qagredo_workflow_part3_grading_output_16x9.png",
}

def add_image_native(slide, path: Path) -> None:
    if path.exists():
        # Insert at native image size (no explicit width/height scaling).
        # Images are pre-rendered at 16:9 to match slide dimensions.
        slide.shapes.add_picture(str(path), Inches(0), Inches(0))


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: overview (native full-slide image)
    s = prs.slides.add_slide(blank)
    add_image_native(s, IMG["full"])

    # Slide 2: part 1 (native full-slide image)
    s = prs.slides.add_slide(blank)
    add_image_native(s, IMG["part1"])

    # Slide 3: part 2 (native full-slide image)
    s = prs.slides.add_slide(blank)
    add_image_native(s, IMG["part2"])

    # Slide 4: part 3 (native full-slide image)
    s = prs.slides.add_slide(blank)
    add_image_native(s, IMG["part3"])

    # Slide 5: framework upgrade summary
    s = prs.slides.add_slide(blank)
    title = s.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8.5), Inches(0.8))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Framework Upgrade Outcome"
    p.font.size = Pt(34)
    p.font.bold = True
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.8), Inches(4.8))
    btf = body.text_frame
    btf.clear()
    for i, line in enumerate(
        [
            "LangChain: prompt templates + structured parsing inside generation nodes.",
            "LangGraph: per-document orchestration with state, routing, and fallback.",
            "Input conversion remains parser-based; format support is unchanged.",
            "2-GPU deployment model remains unchanged unless extra model calls are added.",
        ]
    ):
        bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        bp.text = line
        bp.font.size = Pt(20)

    OUTFILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTFILE))
    print(str(OUTFILE))


if __name__ == "__main__":
    build()
