#!/usr/bin/env python3
"""
Build editable PPTX for QAGRedo workflow (current per-document flow).

Native shapes + connectors — editable in PowerPoint.
Mirrors docs/architecture/diagrams/qagredo_workflow_current.dot.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

_FILL = RGBColor(248, 250, 252)
_LINE = RGBColor(71, 85, 105)
_TEXT = RGBColor(15, 23, 42)
_EDGE = RGBColor(29, 78, 216)
_EDGE_FAIL = RGBColor(239, 68, 68)
_ACCEPT_FILL = RGBColor(220, 252, 231)
_ACCEPT_LINE = RGBColor(22, 163, 74)


def _box(
    slide,
    left,
    top,
    w,
    h,
    text: str,
    *,
    fill_rgb: RGBColor | None = None,
    line_rgb: RGBColor | None = None,
) -> object:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_rgb or _FILL
    sh.line.color.rgb = line_rgb or _LINE
    sh.line.width = Pt(1.1)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.name = "Calibri"
    p.font.color.rgb = _TEXT
    p.alignment = PP_ALIGN.CENTER
    return sh


def _diamond(slide, left, top, side, text: str) -> object:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.FLOWCHART_DECISION, left, top, side, side
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = _FILL
    sh.line.color.rgb = _LINE
    sh.line.width = Pt(1.1)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(7.5)
    p.font.name = "Calibri"
    p.font.color.rgb = _TEXT
    p.alignment = PP_ALIGN.CENTER
    return sh


def _arrow(
    slide,
    x1: Emu,
    y1: Emu,
    x2: Emu,
    y2: Emu,
    *,
    dashed: bool = False,
    label: str | None = None,
    color: RGBColor | None = None,
) -> None:
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, int(x1), int(y1), int(x2), int(y2)
    )
    conn.line.color.rgb = color or _EDGE
    conn.line.width = Pt(1.2)
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if label:
        mx = (int(x1) + int(x2)) // 2
        my = (int(y1) + int(y2)) // 2
        tb = slide.shapes.add_textbox(
            Emu(mx - 457200), Emu(my - 133350), Emu(914400), Emu(266700)
        )
        tb.fill.background()
        tb.line.fill.background()
        p = tb.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(7.5)
        p.font.name = "Calibri"
        p.font.color.rgb = color or _EDGE


def _curve_fail(slide, x1, y1, x2, y2) -> None:
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.CURVE, int(x1), int(y1), int(x2), int(y2)
    )
    conn.line.color.rgb = _EDGE_FAIL
    conn.line.width = Pt(1.75)


def build(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = int(13.3333 * 914400)
    prs.slide_height = int(7.5 * 914400)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.18), Inches(12.6), Inches(0.52)
    )
    title_box.fill.background()
    title_box.line.fill.background()
    tp = title_box.text_frame.paragraphs[0]
    tp.text = (
        "QAGRedo workflow — per-document execution (editable)"
    )
    tp.font.size = Pt(20)
    tp.font.bold = True
    tp.font.name = "Calibri"
    tp.font.color.rgb = _TEXT

    y_main = Inches(1.52)
    w, h = Inches(0.9), Inches(0.66)
    gap = Inches(0.11)
    start_x = Inches(0.28)

    def cx(i: int):
        return start_x + i * (w + gap)

    src = _box(slide, cx(0), y_main, w, h, "Source\ndocument")
    init_q = _box(
        slide, cx(1), y_main, w, h, "Generate initial\nquestions (N)"
    )
    slot_q = _box(
        slide, cx(2), y_main, w, h, "Current question\n(slot i)"
    )
    ans = _box(slide, cx(3), y_main, w, h, "Answer\ngeneration")
    sem = _box(
        slide,
        cx(4),
        y_main,
        w,
        h,
        "Keyword check\n(optional pass)",
    )

    d_side = Inches(0.68)
    gate = _diamond(
        slide,
        cx(5) + Inches(0.11),
        y_main + Inches(0.02),
        d_side,
        "Grounded +\nconf ≥\nthreshold?",
    )

    accept = _box(
        slide,
        cx(6) + Inches(0.06),
        y_main,
        Inches(0.88),
        h,
        "Accept pair",
        fill_rgb=_ACCEPT_FILL,
        line_rgb=_ACCEPT_LINE,
    )

    y_llm = Inches(2.62)
    llm = _box(
        slide,
        cx(4),
        y_llm,
        w,
        Inches(0.72),
        "LLM judge fallback\n(Qwen) if needed",
    )

    rep_y = Inches(3.38)
    regen = _box(
        slide,
        cx(5) - Inches(0.05),
        rep_y,
        Inches(1.15),
        Inches(0.78),
        "Regenerate question\nfor this slot",
    )

    save_txt = (
        "Save the QA pairs\n"
        "(after all N slots)\n\n"
        "• split supporting_evidence: strip list prefixes,\n"
        "  dedupe fragments, then map to citation_spans\n"
        "• citation_notes for fragments not in document\n"
        "• citation_alignment policy "
        "(off / annotate / strict)\n"
        "• hallucination_check in qa_pairs[]\n"
        "• grading_summary from final qa_pairs\n"
        "• save doc_<id>_analysis.json"
    )
    save_blk = _box(
        slide,
        Inches(8.25),
        Inches(1.38),
        Inches(4.75),
        Inches(3.05),
        save_txt,
    )
    tf_s = save_blk.text_frame
    tf_s.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf_s.paragraphs[0].font.size = Pt(9)

    schema_txt = (
        "Final output schema\n"
        "• document\n"
        "• qa_pairs (question, answer,\n"
        "  hallucination_check,\n"
        "  citation_spans, citation_notes)\n"
        "• question_generation\n"
        "• answer_generation\n"
        "• grading_summary\n"
        "• run_metrics"
    )
    note = _box(
        slide,
        Inches(8.25),
        Inches(4.58),
        Inches(4.75),
        Inches(2.55),
        schema_txt,
    )
    note_tf = note.text_frame
    note_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    note_tf.paragraphs[0].font.size = Pt(9)

    # Main chain
    chain = [src, init_q, slot_q, ans, sem]
    for i in range(len(chain) - 1):
        a, b = chain[i], chain[i + 1]
        _arrow(
            slide,
            Emu(a.left + a.width),
            Emu(a.top + a.height // 2),
            Emu(b.left),
            Emu(b.top + b.height // 2),
        )

    _arrow(
        slide,
        Emu(sem.left + sem.width),
        Emu(sem.top + sem.height // 2),
        Emu(gate.left),
        Emu(gate.top + gate.height // 2),
        dashed=True,
        label="clear / high\nconfidence",
    )
    _arrow(
        slide,
        Emu(sem.left + sem.width // 2),
        Emu(sem.top + sem.height),
        Emu(llm.left + llm.width // 2),
        Emu(llm.top),
        label="low conf /\nflagged",
    )
    _arrow(
        slide,
        Emu(llm.left + llm.width),
        Emu(llm.top + llm.height // 2),
        Emu(gate.left + gate.width // 2),
        Emu(gate.top + gate.height),
    )

    _arrow(
        slide,
        Emu(gate.left + gate.width),
        Emu(gate.top + gate.height // 2),
        Emu(accept.left),
        Emu(accept.top + accept.height // 2),
        label="yes",
    )
    _arrow(
        slide,
        Emu(gate.left + gate.width // 2),
        Emu(gate.top + gate.height),
        Emu(regen.left + regen.width // 2),
        Emu(regen.top),
        label="no",
        color=_EDGE_FAIL,
    )
    _curve_fail(
        slide,
        Emu(regen.left),
        Emu(regen.top + regen.height // 2),
        Emu(slot_q.left + slot_q.width // 2),
        Emu(slot_q.top + slot_q.height),
    )
    rtb = slide.shapes.add_textbox(
        Inches(2.85), Inches(3.55), Inches(2.1), Inches(0.4)
    )
    rtb.fill.background()
    rtb.line.fill.background()
    rp = rtb.text_frame.paragraphs[0]
    rp.text = "retry (max_question_regeneration_rounds)"
    rp.font.size = Pt(8)
    rp.font.bold = True
    rp.font.name = "Calibri"
    rp.font.color.rgb = _EDGE_FAIL

    _arrow(
        slide,
        Emu(accept.left + accept.width),
        Emu(accept.top + accept.height // 2),
        Emu(save_blk.left),
        Emu(save_blk.top + save_blk.height // 2),
        dashed=True,
        label="all slots done",
    )

    foot = slide.shapes.add_textbox(
        Inches(0.28), Inches(6.85), Inches(12.75), Inches(0.48)
    )
    foot.fill.background()
    foot.line.fill.background()
    fp = foot.text_frame.paragraphs[0]
    fp.text = (
        "Source: run_qa_pipeline.py + utils/hallucination_checker.py "
        "(hybrid grading). Regenerate: "
        "python3 scripts/utils/build_qagredo_workflow_current_pptx.py"
    )
    fp.font.size = Pt(8.5)
    fp.font.name = "Calibri"
    fp.font.color.rgb = _LINE

    prs.save(str(output))
    print(f"Wrote {output}")


def main() -> None:
    ap = argparse.ArgumentParser(
        description="Build editable workflow-current PPTX."
    )
    default_out = Path(
        "docs/architecture/diagrams/"
        "qagredo_workflow_current_editable.pptx"
    )
    ap.add_argument(
        "-o",
        "--output",
        type=Path,
        default=default_out,
        help="Output .pptx path",
    )
    args = ap.parse_args()
    build(args.output.resolve())


if __name__ == "__main__":
    main()
