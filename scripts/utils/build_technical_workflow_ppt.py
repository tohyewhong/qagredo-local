from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTFILE = Path(
    "/home/tyewhong/qagredo/docs/QAGRedo_Technical_Workflow_20slides.pptx"
)
ASSET_DIR = Path("/home/tyewhong/.cursor/projects/home-tyewhong/assets")
DOCS_DIR = Path("/home/tyewhong/qagredo/docs")

IMG = {
    "hero": ASSET_DIR / "slide_visual_01_title_hero.png",
    "problem": ASSET_DIR / "slide_visual_02_problem_context.png",
    "pipeline": DOCS_DIR / "qagredo_full_pipeline_flow_16x9.png",
    "arch": ASSET_DIR / "slide_visual_04_architecture_containers.png",
    "qgen": ASSET_DIR / "slide_visual_05_question_generation.png",
    "aval": ASSET_DIR / "slide_visual_06_answer_validation.png",
    "hybrid": ASSET_DIR / "slide_visual_07_hybrid_grading.png",
    "audit": ASSET_DIR / "slide_visual_08_output_audit.png",
    "config": ASSET_DIR / "slide_visual_09_config_tuning.png",
    "tradeoff": ASSET_DIR / "slide_visual_10_tradeoff_balance.png",
    "tradeoff_explicit": ASSET_DIR
    / "slide_visual_16_tradeoff_matrix_explicit.png",
    "offline": ASSET_DIR / "slide_visual_11_offline_deployment.png",
    "risk": ASSET_DIR / "slide_visual_12_reliability_risk.png",
    "qalgo": ASSET_DIR / "slide_visual_13_question_algorithm_flow.png",
    "aprompt": ASSET_DIR / "slide_visual_14_answer_prompt_diagram.png",
}


def add_title_box(slide, title, subtitle=""):
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.9)
    )
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(18, 47, 95)

    if subtitle:
        sub = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.5)
        )
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(15)
        sp.font.color.rgb = RGBColor(66, 82, 108)


def add_bullets(slide, items, x=0.7, y=1.7, w=6.2, h=4.8, font_size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(35, 48, 69)
        p.space_after = Pt(8)


def add_image(slide, path, x, y, w, h):
    if path.exists():
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.color.rgb = RGBColor(220, 228, 240)
        slide.shapes.add_picture(
            str(path), Inches(x), Inches(y), Inches(w), Inches(h)
        )


def add_code_block(slide, title, code, x=0.7, y=1.7, w=12.0, h=4.8):
    add_title_box(slide, title)
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(17, 27, 45)
    bg.line.color.rgb = RGBColor(17, 27, 45)
    tf = bg.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = "Consolas"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(220, 230, 245)


def add_table_like(slide, headers, rows, x=0.7, y=1.8, w=12.0, h=4.8):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 248, 255)
    shape.line.color.rgb = RGBColor(170, 185, 210)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()

    hp = tf.paragraphs[0]
    hp.text = " | ".join(headers)
    hp.font.bold = True
    hp.font.size = Pt(15)
    hp.font.color.rgb = RGBColor(26, 56, 106)
    hp.space_after = Pt(8)

    for row in rows:
        p = tf.add_paragraph()
        p.text = " | ".join(row)
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(46, 56, 76)
        p.space_after = Pt(6)


def add_stage_box(
    slide,
    x,
    y,
    w,
    h,
    title,
    lines,
    fill_rgb,
    title_font_size=14,
    body_font_size=11,
):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*fill_rgb)
    box.line.color.rgb = RGBColor(160, 175, 205)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(title_font_size)
    p.font.color.rgb = RGBColor(25, 45, 85)
    for line in lines:
        bp = tf.add_paragraph()
        bp.text = line
        bp.font.size = Pt(body_font_size)
        bp.font.color.rgb = RGBColor(45, 58, 80)
        bp.space_after = Pt(2)
    return box


def add_arch_label(slide, x, y, w, h, title, detail, fill_rgb):
    lbl = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = RGBColor(*fill_rgb)
    lbl.line.color.rgb = RGBColor(120, 145, 190)
    tf = lbl.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(23, 46, 86)
    d = tf.add_paragraph()
    d.text = detail
    d.font.size = Pt(11)
    d.font.color.rgb = RGBColor(45, 58, 80)
    return lbl


def add_overlay_label(
    slide,
    x,
    y,
    w,
    h,
    text,
    fill_rgb=(255, 255, 255),
    font_rgb=(25, 45, 85),
    font_size=11,
):
    label = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    label.fill.solid()
    label.fill.fore_color.rgb = RGBColor(*fill_rgb)
    label.fill.transparency = 8
    label.line.color.rgb = RGBColor(140, 160, 195)
    tf = label.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*font_rgb)
    p.alignment = PP_ALIGN.CENTER
    return label


def add_takeaway(slide, text):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(6.75),
        Inches(12.0),
        Inches(0.45),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(240, 246, 255)
    box.line.color.rgb = RGBColor(170, 190, 225)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Key takeaway: {text}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 56, 106)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    add_image(s, IMG["hero"], 0, 0, 13.333, 7.5)
    overlay = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
    overlay.fill.transparency = 18
    overlay.line.fill.background()
    t = s.shapes.add_textbox(
        Inches(0.7), Inches(1.6), Inches(10.8), Inches(2.2)
    ).text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = "QAGRedo Workflow + LangChain/LangGraph Upgrade"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(18, 47, 95)
    p = t.add_paragraph()
    p.text = "Minimal Workflow, Algorithms, and Engineering Tradeoffs"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(66, 82, 108)

    # 2
    s = prs.slides.add_slide(blank)
    add_title_box(
        s, "2. Why This Workflow Exists", "Business and technical context"
    )
    add_bullets(
        s,
        [
            "Goal: generate complex Q&A grounded strictly in source documents.",  # noqa: E501
            "Problem: raw LLM output is fast but can be shallow, repetitive, or hallucinatory.",  # noqa: E501
            "Design response: multi-stage validation with adaptive fallback checks.",  # noqa: E501
            "Outcome target: high factual confidence with auditable evidence trail.",  # noqa: E501
        ],
        w=6.5,
        font_size=19,
    )
    add_image(s, IMG["problem"], 7.0, 1.7, 6.0, 4.9)
    # Add on-image concept labels to mirror the 10-slide storytelling style.
    add_overlay_label(
        s,
        7.15,
        5.95,
        1.2,
        0.36,
        "Need",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
        font_size=10,
    )
    add_overlay_label(
        s,
        8.45,
        5.95,
        1.2,
        0.36,
        "Risk",
        fill_rgb=(255, 238, 225),
        font_rgb=(118, 52, 18),
        font_size=10,
    )
    add_overlay_label(
        s,
        9.75,
        5.95,
        1.45,
        0.36,
        "Response",
        fill_rgb=(233, 252, 244),
        font_rgb=(14, 86, 48),
        font_size=10,
    )
    add_overlay_label(
        s,
        11.35,
        5.95,
        1.2,
        0.36,
        "Target",
        fill_rgb=(240, 236, 255),
        font_rgb=(58, 36, 120),
        font_size=10,
    )

    # 3
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "3. Minimal End-to-End Workflow",
        "Full pipeline from input to audited output",
    )
    add_image(s, IMG["pipeline"], 0.6, 1.3, 12.1, 3.35)
    # Keep the full-flow image clean; summarize using lane cards below.
    add_stage_box(
        s,
        0.8,
        4.95,
        3.95,
        1.35,
        "Lane A: Input + Preparation",
        [
            "Input routing by file type",
            "Normalization + optional semantic enrich",
        ],
        (232, 241, 255),
        title_font_size=16,
        body_font_size=12,
    )
    add_stage_box(
        s,
        4.95,
        4.95,
        3.95,
        1.35,
        "Lane B: Q/A Quality Loop",
        [
            "LangChain prompt/template/parsing path",
            "Grounding, retries, and coverage checks",
        ],
        (233, 252, 244),
        title_font_size=16,
        body_font_size=12,
    )
    add_stage_box(
        s,
        9.1,
        4.95,
        3.4,
        1.35,
        "Lane C: Grading + Output",
        [
            "Hybrid scoring path",
            "Grade, issues, run summary + run metrics",
        ],
        (240, 236, 255),
        title_font_size=16,
        body_font_size=12,
    )
    add_takeaway(
        s,
        "Single visual, three lanes: preparation, quality loop, and final audit output.",  # noqa: E501
    )

    # 4
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "4. Runtime Architecture",
        "Three-container design with explicit runtime responsibilities",
    )
    add_image(s, IMG["arch"], 0.6, 1.4, 7.2, 5.5)
    # On-diagram container labels to make architecture self-explanatory
    add_arch_label(
        s,
        0.95,
        5.55,
        2.15,
        1.0,
        "qagredo-vllm",
        "Llama Q/A generator\nGPU 0 | Port ${VLLM_HOST_PORT}",
        (232, 241, 255),
    )
    add_arch_label(
        s,
        3.05,
        5.55,
        2.2,
        1.0,
        "qagredo-runner",
        "Orchestrator\nCPU | internal",
        (233, 252, 244),
    )
    add_arch_label(
        s,
        5.25,
        5.55,
        2.35,
        1.0,
        "qagredo-vllm-judge",
        "Qwen judge model\nGPU 1 | Port ${VLLM_JUDGE_HOST_PORT}",
        (240, 236, 255),
    )

    # High-contrast side cards instead of dense table
    add_stage_box(
        s,
        8.15,
        1.75,
        4.45,
        1.25,
        "Generator service",
        [
            "qagredo-vllm | Llama | GPU 0 | :${VLLM_HOST_PORT}",
            "Handles question + answer generation",
        ],
        (232, 241, 255),
    )
    add_stage_box(
        s,
        8.15,
        3.2,
        4.45,
        1.25,
        "Judge service",
        [
            "qagredo-vllm-judge | Qwen | GPU 1 | :${VLLM_JUDGE_HOST_PORT}",
            "Independent grounding verification",
        ],
        (240, 236, 255),
    )
    add_stage_box(
        s,
        8.15,
        4.65,
        4.45,
        1.25,
        "Runner service",
        [
            "qagredo-runner | CPU | internal network",
            "Pipeline orchestration + semantic checks",
        ],
        (233, 252, 244),
    )
    add_bullets(
        s,
        [
            "Separate judge model removes self-evaluation bias.",
            "GPU split avoids head-of-line blocking between generation and grading.",  # noqa: E501
        ],
        x=8.15,
        y=6.0,
        w=4.45,
        h=1.2,
        font_size=12,
    )

    # 5
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "5. Pipeline Orchestration Flow Diagram",
        "Execution sequence with key gates",
    )
    add_stage_box(
        s,
        0.7,
        1.9,
        2.25,
        2.2,
        "Input",
        [
            "Load supported input",
            "Normalize to JSONL",
            "Optional semantic metadata",
            "Prepare run metadata",
        ],
        (232, 241, 255),
    )
    add_stage_box(
        s,
        3.2,
        1.9,
        2.4,
        2.2,
        "Question Gen",
        [
            "Prompt + few-shot",
            "Dedup (LLM or Jaccard)",
            "Grounding + comp checks",
        ],
        (226, 245, 255),
    )
    add_stage_box(
        s,
        5.9,
        1.9,
        2.4,
        2.2,
        "Answer Gen",
        [
            "Structured answer",
            "Grounding retries",
            "Coverage rewrite pass",
        ],
        (232, 252, 242),
    )
    add_stage_box(
        s,
        8.6,
        1.9,
        2.3,
        2.2,
        "Grading",
        [
            "Hybrid checker",
            "Semantic first",
            "Judge fallback",
        ],
        (237, 233, 254),
    )
    add_stage_box(
        s,
        11.15,
        1.9,
        1.5,
        2.2,
        "Output",
        [
            "analysis.json",
            "run_summary.json",
            "run_metrics",
            "timestamped path",
        ],
        (255, 238, 225),
    )

    # Arrows between boxes
    for x0, x1 in [(2.95, 3.18), (5.62, 5.88), (8.32, 8.58), (10.92, 11.12)]:
        arr = s.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            Inches(x0),
            Inches(2.65),
            Inches(x1 - x0),
            Inches(0.45),
        )
        arr.fill.solid()
        arr.fill.fore_color.rgb = RGBColor(83, 130, 212)
        arr.line.color.rgb = RGBColor(83, 130, 212)

    add_bullets(
        s,
        [
            "Deterministic stage order keeps behavior explainable and reproducible.",  # noqa: E501
            "Each stage exposes explicit quality gates before proceeding downstream.",  # noqa: E501
            "Audit outputs preserve both final verdicts and failure reasons.",
        ],
        y=4.7,
        w=12.0,
        font_size=16,
    )

    # 6
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "6. Question Generation Strategy",
        "10 types, benefits, examples, and why few-shot prevents trivial lookups",  # noqa: E501
    )
    add_stage_box(
        s,
        0.65,
        1.6,
        5.9,
        3.85,
        "Types 1-5 (Core Analytical)",
        [
            "1) Analysis -> uncover drivers and dependencies",
            "2) Aggregation -> improve counting accuracy",
            "3) Comparison -> expose subtle differences",
            "4) Inference -> test implied understanding",
            "5) Causal -> validate cause-effect logic",
        ],
        (233, 242, 255),
    )
    add_stage_box(
        s,
        6.75,
        1.6,
        5.9,
        3.85,
        "Types 6-10 (Advanced Reasoning)",
        [
            "6) Temporal -> verify timeline consistency",
            "7) Multi-hop -> connect distant facts",
            "8) Synthesis -> integrate 3+ evidence points",
            "9) Evaluation -> rate claim support strength",
            "10) Counterfactual -> test causal robustness",
        ],
        (233, 252, 244),
    )
    add_bullets(
        s,
        [
            "Few-shot effect: GOOD and BAD examples teach the model what 'complex' looks like, reducing trivial lookup questions.",  # noqa: E501
            "Practical impact: fewer weak questions pass forward, so downstream answer retries are lower.",  # noqa: E501
        ],
        x=0.7,
        y=5.65,
        w=12.0,
        h=1.0,
        font_size=13,
    )

    # 7
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "7. Question Generation Algorithm",
        "Pictorial flow: generate, filter, validate, and loop",
    )
    add_image(s, IMG["qalgo"], 0.7, 1.45, 12.0, 3.7)
    # Add in-image labels so the flow is understandable at a glance.
    add_overlay_label(s, 0.9, 2.95, 1.9, 0.38, "1) Prompt + LLM")
    add_overlay_label(s, 4.1, 2.95, 1.8, 0.38, "2) Parse")
    add_overlay_label(s, 5.95, 2.35, 2.15, 0.38, "3) Dedup")
    add_overlay_label(s, 5.95, 2.95, 2.15, 0.38, "4) Quality Gates")
    add_overlay_label(s, 8.35, 2.35, 2.2, 0.38, "5) Accept Path")
    add_overlay_label(s, 8.35, 2.95, 2.2, 0.38, "6) Regenerate Loop")
    add_overlay_label(s, 10.85, 2.65, 1.55, 0.38, "7) Final Set")
    add_bullets(
        s,
        [
            "Flow interpretation: prompt build -> LLM generation -> parse -> dedup -> quality gates -> final validated set.",  # noqa: E501
            "Upper path enforces grounding/comprehensiveness acceptance; lower loop represents targeted regeneration cycles.",  # noqa: E501
            "Tradeoff: extra loop cost improves final question diversity and reduces downstream answer failures.",  # noqa: E501
        ],
        y=5.35,
        w=12.1,
        h=1.9,
        font_size=14,
    )

    # 8
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "8. Question Quality Gates",
        "Grounding + comprehensiveness hard gates",
    )
    add_stage_box(
        s,
        0.8,
        1.9,
        3.7,
        3.3,
        "Gate A: Grounding",
        [
            "Purpose: reject out-of-doc questions",
            "Default: min_conf=0.7, max_retry=2",
            "Tradeoff: precision vs latency",
        ],
        (233, 242, 255),
    )
    add_stage_box(
        s,
        4.8,
        1.9,
        3.7,
        3.3,
        "Gate B: Comprehensiveness",
        [
            "Purpose: reject trivial lookups",
            "Default: min_score=0.6, max_attempts=2",
            "Tradeoff: depth vs extra LLM calls",
        ],
        (233, 252, 244),
    )
    add_stage_box(
        s,
        8.8,
        1.9,
        3.7,
        3.3,
        "Method: Semantic-first",
        [
            "Purpose: keep checks fast",
            "Default for question validation path",
            "Risk: misses some nuanced cases",
        ],
        (240, 236, 255),
    )
    add_bullets(
        s,
        [
            "Net effect: questions are both answerable and meaningful before answer generation starts."  # noqa: E501
        ],
        y=5.55,
        w=12.0,
        h=0.8,
        font_size=16,
    )

    # 9
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "9. Structured Answer Prompt Design",
        "Pictorial prompt constraints and output branches",
    )
    add_image(s, IMG["aprompt"], 0.7, 1.4, 12.0, 3.8)
    # High-contrast labels on dark diagram.
    add_overlay_label(
        s,
        0.85,
        2.15,
        2.2,
        0.38,
        "Document + Question",
        fill_rgb=(225, 243, 255),
        font_rgb=(16, 40, 85),
    )
    add_overlay_label(
        s,
        3.35,
        2.15,
        2.2,
        0.38,
        "Prompt Constraints",
        fill_rgb=(225, 243, 255),
        font_rgb=(16, 40, 85),
    )
    add_overlay_label(
        s,
        5.75,
        2.15,
        1.95,
        0.38,
        "Answer Gen",
        fill_rgb=(225, 243, 255),
        font_rgb=(16, 40, 85),
    )
    add_overlay_label(
        s,
        8.0,
        1.8,
        2.35,
        0.38,
        "Counting-First Branch",
        fill_rgb=(225, 255, 238),
        font_rgb=(14, 86, 48),
    )
    add_overlay_label(
        s,
        8.0,
        2.65,
        2.35,
        0.38,
        "Insufficient-Info Branch",
        fill_rgb=(255, 238, 225),
        font_rgb=(118, 52, 18),
    )
    add_overlay_label(
        s,
        10.55,
        2.2,
        2.05,
        0.38,
        "Answer + Evidence",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
    )
    add_overlay_label(
        s,
        7.95,
        3.05,
        2.45,
        0.38,
        "Insufficient Info Signal",
        fill_rgb=(255, 238, 225),
        font_rgb=(118, 52, 18),
    )
    add_overlay_label(
        s,
        10.55,
        2.65,
        2.05,
        0.38,
        "Fallback Response",
        fill_rgb=(255, 238, 225),
        font_rgb=(118, 52, 18),
    )
    add_bullets(
        s,
        [
            "The visual maps prompt constraints to output behavior: answer + supporting evidence or explicit insufficiency path.",  # noqa: E501
            "Counting-first branch enforces itemization before totals, reducing aggregation mistakes.",  # noqa: E501
            "Branch-based prompt design increases factual consistency at temperature 0.3 while preserving auditability.",  # noqa: E501
        ],
        y=5.35,
        w=12.0,
        h=1.9,
        font_size=14,
    )
    add_takeaway(
        s, "Prompt structure enforces both factuality and explainability."
    )

    # 10
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "10. Answer Validation and Retry Loop",
        "Grounding-first acceptance policy",
    )
    add_image(s, IMG["aval"], 6.6, 1.4, 6.5, 5.5)
    # Add in-picture words so the validation flow is self-explanatory.
    add_overlay_label(
        s,
        6.75,
        1.75,
        2.1,
        0.38,
        "Generated Answer",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
    )
    add_overlay_label(
        s,
        8.95,
        3.25,
        1.7,
        0.38,
        "Grounding Check",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
    )
    add_overlay_label(
        s,
        10.75,
        1.95,
        2.2,
        0.38,
        "Supported Evidence",
        fill_rgb=(225, 255, 238),
        font_rgb=(14, 86, 48),
    )
    add_overlay_label(
        s,
        10.75,
        3.65,
        2.2,
        0.38,
        "Ungrounded Claims",
        fill_rgb=(255, 238, 225),
        font_rgb=(118, 52, 18),
    )
    add_overlay_label(
        s,
        9.05,
        5.25,
        2.15,
        0.38,
        "Accept / Retry",
        fill_rgb=(255, 247, 220),
        font_rgb=(96, 73, 15),
    )
    add_bullets(
        s,
        [
            "Initial answer is checked by hallucination checker.",
            "If not grounded at required confidence -> regenerate up to 3 times.",  # noqa: E501
            "Best available answer is retained for transparency, not silently dropped.",  # noqa: E501
        ],
        w=6.0,
        font_size=16,
    )

    # 11
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "11. Coverage Validator + Targeted Rewrite",
        "New quality layer for question-answer alignment",
    )
    add_image(s, IMG["aval"], 6.7, 1.6, 6.0, 5.0)
    # In-picture labels to explain the coverage-rewrite flow at a glance.
    add_overlay_label(
        s,
        6.85,
        1.95,
        1.95,
        0.38,
        "Initial Answer",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
    )
    add_overlay_label(
        s,
        8.9,
        3.0,
        1.9,
        0.38,
        "Coverage Check",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
    )
    add_overlay_label(
        s,
        10.9,
        1.95,
        1.75,
        0.38,
        "Covered",
        fill_rgb=(225, 255, 238),
        font_rgb=(14, 86, 48),
    )
    add_overlay_label(
        s,
        10.9,
        3.55,
        1.75,
        0.38,
        "Low Coverage",
        fill_rgb=(255, 238, 225),
        font_rgb=(118, 52, 18),
    )
    add_overlay_label(
        s,
        8.75,
        5.05,
        2.2,
        0.38,
        "Targeted Rewrite (x1)",
        fill_rgb=(255, 247, 220),
        font_rgb=(96, 73, 15),
    )
    add_overlay_label(
        s,
        11.05,
        5.05,
        1.65,
        0.38,
        "Grounding Gate",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
        font_size=10,
    )
    add_bullets(
        s,
        [
            "Purpose: detect grounded-but-incomplete answers that miss parts of the question.",  # noqa: E501
            "If coverage is low, run one targeted rewrite using missing-point feedback.",  # noqa: E501
            "Safety: rewritten answer is accepted only if grounding/confidence still passes.",  # noqa: E501
            "Outcome: improves completeness without increasing hallucination risk.",  # noqa: E501
        ],
        x=0.7,
        y=1.8,
        w=5.7,
        h=4.8,
        font_size=15,
    )

    # 12
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "12. Hallucination Checking Methods",
        "Method-level strengths and limits",
    )
    add_stage_box(
        s,
        0.8,
        1.9,
        3.7,
        3.6,
        "Semantic",
        [
            "Strength: fast, CPU, no GPU",
            "Weakness: limited for counting/inference",
            "Use: first-pass filter",
        ],
        (233, 242, 255),
    )
    add_stage_box(
        s,
        4.8,
        1.9,
        3.7,
        3.6,
        "LLM Judge",
        [
            "Strength: handles reasoning/aggregation",
            "Weakness: slower and GPU-expensive",
            "Use: high-ambiguity cases",
        ],
        (240, 236, 255),
    )
    add_stage_box(
        s,
        8.8,
        1.9,
        3.7,
        3.6,
        "Hybrid (Default)",
        [
            "Strength: speed-accuracy balance",
            "Weakness: extra routing logic",
            "Use: production default path",
        ],
        (233, 252, 244),
    )
    add_bullets(
        s,
        [
            "Hybrid captures most cases cheaply, escalates only uncertain edges.",  # noqa: E501
            "Independent judge model gives stronger quality signal than self-judging.",  # noqa: E501
        ],
        y=5.75,
        w=12.0,
        font_size=14,
    )

    # 13
    s = prs.slides.add_slide(blank)
    add_title_box(
        s, "13. Hybrid Grading Decision Path", "Adaptive routing logic"
    )
    add_image(s, IMG["hybrid"], 0.7, 1.4, 6.4, 5.2)
    # In-image labels to link visual flow to explanation bullets.
    add_overlay_label(
        s,
        1.0,
        1.8,
        2.15,
        0.38,
        "A) Fast Semantic Path",
        fill_rgb=(225, 255, 238),
        font_rgb=(14, 86, 48),
        font_size=10,
    )
    add_overlay_label(
        s,
        1.0,
        2.65,
        2.15,
        0.38,
        "B) Judge Fallback Path",
        fill_rgb=(240, 236, 255),
        font_rgb=(58, 36, 120),
        font_size=10,
    )
    add_overlay_label(
        s,
        4.5,
        3.1,
        2.25,
        0.38,
        "C) Merge + Final Verdict",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
        font_size=10,
    )
    add_bullets(
        s,
        [
            "A) Fast Semantic Path: high-confidence answers are accepted quickly.",  # noqa: E501
            "B) Judge Fallback Path: low-confidence cases invoke the independent LLM judge.",  # noqa: E501
            "C) Merge + Final Verdict: confidence combines evidence from both checks.",  # noqa: E501
            "Tradeoff: slight routing complexity yields major efficiency gains.",  # noqa: E501
        ],
        x=7.3,
        y=1.8,
        w=5.6,
        h=4.8,
        font_size=15,
    )

    # 14
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "14. Output Schema Made Simple",
        "How to read one result file in under 30 seconds",
    )
    add_image(s, IMG["audit"], 8.0, 1.5, 4.7, 2.5)
    add_stage_box(
        s,
        0.7,
        1.8,
        3.8,
        2.1,
        "Block 1: Document context",
        [
            "document.id, title, source",
            "Input identity for traceability",
        ],
        (233, 242, 255),
    )
    add_stage_box(
        s,
        4.85,
        1.8,
        3.8,
        2.1,
        "Block 2: QA details",
        [
            "qa_pairs: question + answer + grading",
            "Per-pair confidence and issues",
        ],
        (233, 252, 244),
    )
    add_stage_box(
        s,
        0.7,
        4.2,
        3.8,
        2.1,
        "Block 3: Generation metadata",
        [
            "question_generation / answer_generation",
            "Models, provider, timestamps",
        ],
        (240, 236, 255),
    )
    add_stage_box(
        s,
        4.85,
        4.2,
        3.8,
        2.1,
        "Block 4: Summary + enrichment",
        [
            "grading_summary + run_metrics",
            "semantic_enrichment (optional)",
        ],
        (255, 240, 229),
    )
    add_bullets(
        s,
        [
            "Interpretation example: if grade is C with low confidence, inspect qa_pairs.issues first, then trace back to question_generation/comprehensiveness metadata.",  # noqa: E501
        ],
        x=8.0,
        y=4.3,
        w=4.9,
        h=1.7,
        font_size=12,
    )
    add_takeaway(
        s,
        "Read outputs as four blocks: context, QA details, generation metadata, final summary.",  # noqa: E501
    )

    # 15
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "15. Configuration Knobs That Matter",
        "Operational tuning from config/config.<profile>.yaml",
    )
    add_image(s, IMG["config"], 6.7, 1.5, 6.2, 5.3)
    # On-image labels make slider visual directly explainable.
    add_overlay_label(
        s,
        7.05,
        5.95,
        1.95,
        0.38,
        "Quality Knob",
        fill_rgb=(225, 255, 238),
        font_rgb=(14, 86, 48),
        font_size=10,
    )
    add_overlay_label(
        s,
        9.2,
        5.95,
        1.95,
        0.38,
        "Latency Knob",
        fill_rgb=(255, 247, 220),
        font_rgb=(96, 73, 15),
        font_size=10,
    )
    add_overlay_label(
        s,
        11.35,
        5.95,
        1.35,
        0.38,
        "Cost Knob",
        fill_rgb=(240, 236, 255),
        font_rgb=(58, 36, 120),
        font_size=10,
    )
    add_stage_box(
        s,
        0.7,
        1.8,
        5.7,
        1.4,
        "Question generation knobs",
        [
            "complexity, num_questions, dedup threshold,",
            "grounding + comprehensiveness thresholds",
        ],
        (233, 242, 255),
    )
    add_stage_box(
        s,
        0.7,
        3.35,
        5.7,
        1.4,
        "Answer generation knobs",
        ["temperature, retry limits,", "coverage validation thresholds"],
        (233, 252, 244),
    )
    add_stage_box(
        s,
        0.7,
        4.9,
        5.7,
        1.4,
        "Grading knobs",
        [
            "hallucination.method and judge model",
            "semantic | keyword | llm | hybrid",
        ],
        (240, 236, 255),
    )
    add_bullets(
        s,
        [
            "Image mapping: Quality Knob -> complexity/coverage strictness.",
            "Latency Knob -> retry/fallback frequency and path depth.",
            "Cost Knob -> judge-model usage and compute footprint.",
        ],
        x=6.9,
        y=1.8,
        w=5.8,
        h=1.1,
        font_size=11,
    )

    # 16
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "16. Whole Workflow Sequence Diagram",
        "From input document to graded output",
    )
    add_stage_box(
        s,
        0.7,
        2.0,
        2.2,
        1.6,
        "Step 1",
        ["Load supported input", "Normalize + enrich metadata"],
        (232, 241, 255),
    )
    add_stage_box(
        s,
        3.2,
        2.0,
        2.2,
        1.6,
        "Step 2",
        ["Generate questions", "Dedup + validate"],
        (226, 245, 255),
    )
    add_stage_box(
        s,
        5.7,
        2.0,
        2.2,
        1.6,
        "Step 3",
        ["Generate answers", "Retry + coverage"],
        (232, 252, 242),
    )
    add_stage_box(
        s,
        8.2,
        2.0,
        2.2,
        1.6,
        "Step 4",
        ["Hybrid grading", "Confidence + issues"],
        (237, 233, 254),
    )
    add_stage_box(
        s,
        10.7,
        2.0,
        2.0,
        1.6,
        "Step 5",
        ["Save outputs", "Emit summary"],
        (255, 238, 225),
    )
    for x0, x1 in [(2.95, 3.15), (5.45, 5.65), (7.95, 8.15), (10.45, 10.65)]:
        arr = s.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            Inches(x0),
            Inches(2.6),
            Inches(x1 - x0),
            Inches(0.4),
        )
        arr.fill.solid()
        arr.fill.fore_color.rgb = RGBColor(83, 130, 212)
        arr.line.color.rgb = RGBColor(83, 130, 212)
    add_bullets(
        s,
        [
            "Decision Gate A (QG): reject ungrounded or trivial questions before answer stage.",  # noqa: E501
            "Decision Gate B (AG): retry ungrounded answers; if still failing, generate replacement questions in later rounds (max 3), then apply one coverage rewrite pass.",  # noqa: E501
            "Decision Gate C (Grading): fallback to LLM judge only for low-confidence semantic results.",  # noqa: E501
        ],
        y=4.4,
        w=12.0,
        h=2.2,
        font_size=16,
    )

    # 17
    s = prs.slides.add_slide(blank)
    add_title_box(
        s, "17. Core Tradeoff Matrix", "Accuracy vs latency vs compute cost"
    )
    add_image(s, IMG["tradeoff"], 0.8, 1.5, 5.8, 5.0)
    add_table_like(
        s,
        ["Design Choice", "Upside", "Downside"],
        [
            [
                "High question complexity",
                "Better stress-test of comprehension",
                "More generation instability risk",
            ],
            [
                "Low answer temperature",
                "Factual consistency",
                "Less linguistic flexibility",
            ],
            [
                "Hybrid grading",
                "Better speed/quality balance",
                "More decision logic",
            ],
            [
                "Coverage rewrite pass",
                "Completeness gains",
                "Additional LLM call",
            ],
        ],
        x=6.9,
        y=1.8,
        w=5.8,
        h=4.9,
    )

    # 18
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "18. Reliability and Failure Handling",
        "How the workflow degrades safely",
    )
    add_image(s, IMG["risk"], 6.8, 1.5, 6.1, 5.2)
    # In-picture labels aligned with bullet points for quick presenter mapping.
    add_overlay_label(
        s,
        7.0,
        4.9,
        1.75,
        0.38,
        "A) Input Guard",
        fill_rgb=(233, 246, 255),
        font_rgb=(16, 40, 85),
        font_size=10,
    )
    add_overlay_label(
        s,
        8.55,
        4.9,
        2.0,
        0.38,
        "B) Retry Control",
        fill_rgb=(255, 247, 220),
        font_rgb=(96, 73, 15),
        font_size=10,
    )
    add_overlay_label(
        s,
        10.4,
        4.9,
        2.05,
        0.38,
        "C) Fallback Path",
        fill_rgb=(240, 236, 255),
        font_rgb=(58, 36, 120),
        font_size=10,
    )
    add_overlay_label(
        s,
        9.2,
        2.35,
        2.2,
        0.38,
        "D) Audit Trail",
        fill_rgb=(225, 255, 238),
        font_rgb=(14, 86, 48),
        font_size=10,
    )
    add_bullets(
        s,
        [
            "A) Input Guard: invalid or empty questions are explicitly marked, never silently ignored.",  # noqa: E501
            "B) Retry Control: insufficient-information path prevents forced hallucinated answers.",  # noqa: E501
            "C) Fallback Path: semantic-to-judge escalation handles hard reasoning cases.",  # noqa: E501
            "D) Audit Trail: all issues are stored with reasons for operator review.",  # noqa: E501
        ],
        w=6.0,
        font_size=15,
    )
    add_takeaway(
        s,
        "Map each reliability bullet to its A/B/C/D callout for quick explanation.",  # noqa: E501
    )

    # 19
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "19. Deployment Model",
        "Offline/air-gapped 6-file split-model workflow",
    )
    add_image(s, IMG["offline"], 0.7, 1.5, 6.3, 5.2)
    add_bullets(
        s,
        [
            "6-file split-model transfer: vllm rootfs + qagredo image + llama model + qwen model + embed model + bundle.",  # noqa: E501
            "Model split archives: models_llama.tar.gz (Llama HF tree, typical vLLM-judge) and models_qwen.tar.gz (Qwen HF tree, typical vLLM generator, e.g. Qwen3.5-9B).",  # noqa: E501
            "setup_offline.sh prepares links, loads containers, and validates runtime.",  # noqa: E501
            "run.sh executes pipeline with reproducible directory structure and outputs.",  # noqa: E501
            "Partial runs: bash run.sh -- --resume skips docs with existing *_analysis.json.",  # noqa: E501
            "Design supports regulated environments with no external network dependency and incremental bundle updates.",  # noqa: E501
        ],
        x=7.2,
        y=1.8,
        w=5.6,
        h=4.8,
        font_size=15,
    )

    # 20
    s = prs.slides.add_slide(blank)
    add_title_box(
        s,
        "20. Senior-Level Recommendations",
        "Action plan for next optimization cycle",
    )
    add_bullets(
        s,
        [
            "1) Add per-answer metadata: coverage_score, rewrite_applied, fallback_path.",  # noqa: E501
            "2) Build calibration dashboard: confidence distribution vs manual audit outcomes.",  # noqa: E501
            "3) Introduce workload profiles (fast/standard/strict) via config presets.",  # noqa: E501
            "4) Expand reliability testing with synthetic hard cases for counting and counterfactuals.",  # noqa: E501
            "5) Keep current architecture: strongest balance of explainability and quality.",  # noqa: E501
        ],
        x=0.8,
        y=1.7,
        w=11.9,
        h=3.8,
        font_size=18,
    )
    note = s.shapes.add_textbox(
        Inches(0.8), Inches(5.8), Inches(12.0), Inches(1.0)
    ).text_frame
    note.clear()
    p = note.paragraphs[0]
    p.text = "Visual credits: custom-generated professional illustrations for technical presentation use."  # noqa: E501
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(90, 104, 128)
    p.alignment = PP_ALIGN.CENTER

    OUTFILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTFILE))
    print(str(OUTFILE))


if __name__ == "__main__":
    build()
