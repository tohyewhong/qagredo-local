from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTFILE = Path("/home/tyewhong/qagredo/docs/QAGRedo_Executive_Overview_Tradeoffs.pptx")
DOCS_DIR = Path("/home/tyewhong/qagredo/docs")

IMG = {
    "pipeline": DOCS_DIR / "qagredo_full_pipeline_flow_16x9.png",
    "input": DOCS_DIR / "qagredo_input_prep_explained_16x9.png",
    "langgraph": DOCS_DIR / "qagredo_langchain_detail_16x9.png",
    "langchain_q_lane": DOCS_DIR / "qagredo_langchain_question_lane_16x9.png",
    "langchain_a_lane": DOCS_DIR / "qagredo_langchain_answer_lane_16x9.png",
    "langgraph_detailed": DOCS_DIR / "qagredo_langgraph_pipeline_detailed.png",
    "answer_detailed": DOCS_DIR / "qagredo_answer_generator_detailed.png",
    "part2": DOCS_DIR / "qagredo_workflow_part2_qa_loop_16x9.png",
}


def add_title(slide, title, subtitle=""):
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.85)).text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 44, 96)
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.5)).text_frame
        st.clear()
        sp = st.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(15)
        sp.font.color.rgb = RGBColor(68, 82, 106)


def add_bullets(slide, items, x=0.7, y=1.7, w=6.0, h=4.8, font_size=18):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(35, 48, 69)
        p.space_after = Pt(8)


def add_image(slide, path, x, y, w, h):
    if path.exists():
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.color.rgb = RGBColor(220, 228, 240)
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def add_image_crop(slide, path, x, y, w, h, crop_left=0.0, crop_right=0.0, crop_top=0.0, crop_bottom=0.0):
    if path.exists():
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.color.rgb = RGBColor(220, 228, 240)
        pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
        pic.crop_left = crop_left
        pic.crop_right = crop_right
        pic.crop_top = crop_top
        pic.crop_bottom = crop_bottom


def add_card(slide, title, body_lines, x, y, w, h, fill_rgb):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*fill_rgb)
    box.line.color.rgb = RGBColor(170, 186, 214)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(17, 44, 90)
    for line in body_lines:
        bp = tf.add_paragraph()
        bp.text = line
        bp.font.size = Pt(12)
        bp.font.color.rgb = RGBColor(42, 56, 80)
        bp.space_after = Pt(2)


def add_takeaway(slide, text):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(6.75), Inches(12.0), Inches(0.45))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(238, 245, 255)
    box.line.color.rgb = RGBColor(170, 191, 225)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Executive takeaway: {text}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(24, 55, 104)


def add_source_ref(slide, text):
    ref = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(12.0), Inches(0.3)).text_frame
    ref.clear()
    p = ref.paragraphs[0]
    p.text = f"Source: {text}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(90, 104, 128)
    p.alignment = PP_ALIGN.RIGHT


def add_tradeoff_table(slide):
    rows = 5
    cols = 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.9), Inches(3.9))
    table = table_shape.table
    table.columns[0].width = Inches(2.3)
    table.columns[1].width = Inches(3.3)
    table.columns[2].width = Inches(3.3)
    table.columns[3].width = Inches(3.0)

    headers = ["Decision lever", "Benefit", "Trade-off", "When to choose"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(33, 83, 160)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)

    data = [
        (
            "High question complexity",
            "Stronger comprehension stress-test",
            "Higher generation volatility and tuning effort",
            "Policy/risk analysis, due diligence",
        ),
        (
            "Strict grounding thresholds",
            "Lower hallucination risk",
            "More retries, longer runtime",
            "Regulated or audit-heavy workflows",
        ),
        (
            "Strict LLM judge (default)",
            "Consistent judge-only governance",
            "Higher latency than lightweight grading paths",
            "Quality-critical production deployment",
        ),
        (
            "Coverage rewrite enabled",
            "More complete answers",
            "Additional model call in some cases",
            "Complex multi-part questions",
        ),
    ]
    for r, row in enumerate(data, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(246, 250, 255)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(36, 50, 74)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title (white background theme)
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    tf = s.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(11.8), Inches(2.2)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "QAGRedo: Executive Overview"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 44, 96)
    p = tf.add_paragraph()
    p.text = "What it does, why it is strong, and where the trade-offs are"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(68, 82, 106)

    # 2. What is QAGRedo
    s = prs.slides.add_slide(blank)
    add_title(s, "1. What QAGRedo Does", "Offline document QA with grounding verification")
    add_bullets(
        s,
        [
            "Ingests heterogeneous documents and normalizes to a canonical JSONL format.",
            "Generates complex questions (not only fact lookups) from each document.",
            "Generates answers with evidence and validates grounding quality before final output.",
            "Grades each Q&A pair and emits auditable artifacts (issues, confidence, traceable metadata).",
        ],
        x=0.7,
        y=1.8,
        w=6.1,
        h=4.8,
        font_size=16,
    )
    add_image(s, IMG["pipeline"], 6.95, 1.7, 5.9, 4.7)
    add_takeaway(s, "QAGRedo is an auditable QA evaluation pipeline, not just a single prompt call.")

    # 3. Why this architecture
    s = prs.slides.add_slide(blank)
    add_title(s, "2. Why This Architecture Is Good", "Separation of concerns improves trust and operability")
    add_card(s, "Generator model", ["Produces questions and answers", "Optimized for useful content generation"], 0.8, 1.8, 3.9, 2.0, (232, 241, 255))
    add_card(s, "Judge model", ["Independent hallucination evaluation", "Reduces self-evaluation bias"], 4.9, 1.8, 3.9, 2.0, (240, 236, 255))
    add_card(s, "Runner orchestration layer", ["Orchestrates stages and quality gates", "Runs strict LLM judge preflight checks"], 9.0, 1.8, 3.4, 2.0, (233, 252, 244))
    add_bullets(
        s,
        [
            "Detailed LangChain parsing flow is shown on the next two zoom slides for readability.",
            "This slide keeps only architecture responsibilities at executive level.",
        ],
        x=0.9,
        y=4.25,
        w=11.6,
        h=1.2,
        font_size=14,
    )
    add_takeaway(s, "Independent judging + staged orchestration increases reliability and explainability.")

    # 3a. LangChain detail zoom (question lane)
    s = prs.slides.add_slide(blank)
    add_title(s, "2A. LangChain Detail Zoom: Question Cleanup Lane", "Zoomed view to improve readability for print")
    add_image(s, IMG["langchain_q_lane"], 0.7, 1.5, 12.0, 5.5)
    add_source_ref(s, "docs/qagredo_langchain_question_lane.dot")
    add_takeaway(s, "This lane shows question text cleanup and normalization before validation.")

    # 3b. LangChain detail zoom (answer lane)
    s = prs.slides.add_slide(blank)
    add_title(s, "2B. LangChain Detail Zoom: Answer Parsing Lane", "Zoomed view to improve readability for print")
    add_image(s, IMG["langchain_a_lane"], 0.7, 1.5, 12.0, 5.5)
    add_source_ref(s, "docs/qagredo_langchain_answer_lane.dot")
    add_takeaway(s, "This lane shows strict JSON parse, extraction fallback, and final structured answer object.")

    # 4. Why LangChain + LangGraph (explicit algorithm rationale)
    s = prs.slides.add_slide(blank)
    add_title(s, "3. Why LangChain and LangGraph Are Used", "Two different problems: step quality vs workflow control")
    add_card(
        s,
        "LangChain (inside each stage)",
        [
            "Purpose: make each LLM call structured and predictable",
            "Used for: prompt templates, output parsing, schema discipline",
            "Benefit: fewer malformed outputs, easier downstream validation",
            "Trade-off: extra abstraction and prompt-management overhead",
        ],
        0.8,
        1.8,
        5.9,
        3.9,
        (232, 241, 255),
    )
    add_card(
        s,
        "LangGraph (across stages)",
        [
            "Purpose: coordinate multi-step pipeline with state and branching",
            "Used for: retries, fallback routing, gate-based progression",
            "Benefit: explicit, auditable control flow for quality governance",
            "Trade-off: more orchestration logic than a simple chain",
        ],
        7.0,
        1.8,
        5.4,
        3.9,
        (233, 252, 244),
    )
    add_bullets(
        s,
        [
            "Professional framing: LangChain optimizes local step reliability; LangGraph optimizes global process reliability.",
            "QAGRedo needs both because quality failures can occur at single-call level and at workflow-routing level.",
        ],
        x=0.9,
        y=5.9,
        w=11.7,
        h=0.8,
        font_size=13,
    )
    add_takeaway(s, "LangChain is the per-step toolbox; LangGraph is the end-to-end traffic controller.")

    # 5. Deep technical flow (node-level)
    s = prs.slides.add_slide(blank)
    add_title(s, "4. Deep Technical Flow: LangGraph Node Routing", "Exact control path from utils/langgraph_pipeline.py")
    add_image(s, IMG["langgraph_detailed"], 0.65, 1.45, 8.2, 5.75)
    add_card(
        s,
        "What this diagram proves",
        [
            "generate_questions -> generate_answers -> grade_primary is the default path",
            "If no questions/answers are produced, flow exits early with skip_reason",
            "grade_fallback_llm is optional compatibility path, not strict default",
            "Question regeneration exists inside question_generator, not as a separate graph node",
        ],
        9.0,
        1.75,
        3.8,
        2.8,
        (245, 248, 255),
    )
    add_card(
        s,
            "Optional fallback trigger (compatibility mode)",
        [
            "enable_dynamic_routing = true",
            "compatibility grading mode is enabled",
            "overall_confidence < fallback_threshold",
        ],
        9.0,
        4.7,
        3.8,
        1.9,
        (255, 240, 229),
    )
    add_source_ref(s, "utils/langgraph_pipeline.py + utils/question_generator.py")
    add_takeaway(s, "Graph-level routing handles stage transitions; question retries happen inside generation internals.")

    # 5a. LangGraph routing zoom
    s = prs.slides.add_slide(blank)
    add_title(s, "4A. LangGraph Routing Zoom", "Core routing and fallback branch (readable print view)")
    add_image_crop(
        s,
        IMG["langgraph_detailed"],
        0.7,
        1.5,
        12.0,
        5.5,
        crop_left=0.12,
        crop_right=0.15,
        crop_top=0.08,
    )
    add_source_ref(s, "docs/qagredo_langgraph_pipeline_detailed.dot")
    add_takeaway(s, "Focus: stage transitions, skip_reason exits, and optional compatibility fallback condition.")

    # 6. Deep technical flow (failed question behavior)
    s = prs.slides.add_slide(blank)
    add_title(s, "5. Failed Question Behavior: Does It Regenerate?", "Yes - via internal validation/comprehensiveness loops")
    add_card(
        s,
        "Grounding validation loop",
        [
            "Function: _validate_and_regenerate_question(...)",
            "Check: is_grounded && confidence >= min_confidence_threshold",
            "If fail: regenerate NEW question and re-check",
            "Retry limit: max_regeneration_attempts (default 2)",
        ],
        0.8,
        1.75,
        6.1,
        2.2,
        (232, 241, 255),
    )
    add_card(
        s,
        "Comprehensiveness loop",
        [
            "Function: _check_question_comprehensiveness(...)",
            "Check: is_comprehensive && score >= comprehensiveness_min_score",
            "If fail: regenerate with weakness-guided prompt",
            "Retry limit: comprehensiveness_max_attempts (default 2)",
        ],
        0.8,
        4.15,
        6.1,
        2.2,
        (233, 252, 244),
    )
    add_card(
        s,
        "Where this runs",
        [
            "Inside generate_questions(...) for each question",
            "Before LangGraph proceeds to answer generation",
            "If all retries fail, best available question is still kept",
        ],
        7.2,
        1.75,
        5.3,
        2.0,
        (240, 236, 255),
    )
    add_card(
        s,
        "Config keys controlling behavior",
        [
            "question_generation.validation.enable_rejection",
            "question_generation.validation.min_confidence_threshold",
            "question_generation.validation.max_regeneration_attempts",
            "question_generation.validation.enable_comprehensiveness_check",
            "question_generation.validation.comprehensiveness_min_score",
            "question_generation.validation.comprehensiveness_max_attempts",
        ],
        7.2,
        3.95,
        5.3,
        2.4,
        (245, 248, 255),
    )
    add_source_ref(s, "utils/question_generator.py")
    add_takeaway(s, "Answer: failed questions are regenerated in two internal loops before moving downstream.")

    # 7. Deep technical flow (answer generator internals)
    s = prs.slides.add_slide(blank)
    add_title(s, "6. Deep Technical Flow: Answer Generator", "Exact logic from utils/answer_generator.py")
    add_image(s, IMG["answer_detailed"], 0.6, 1.45, 8.45, 5.75)
    add_card(
        s,
        "What this diagram proves",
        [
            "Answers are generated, validated, and optionally rewritten before final append",
            "Ungrounded answers trigger regeneration loop up to max_regeneration_attempts",
            "Coverage rewrite is conditional and gated by grounding check",
            "Failure paths are explicit (invalid question / generation failed)",
        ],
        9.2,
        1.75,
        3.5,
        3.0,
        (245, 248, 255),
    )
    add_card(
        s,
        "Acceptance policy",
        [
            "Primary acceptance: grounded + confidence >= threshold",
            "Rewrite accepted only if revised answer is grounded",
            "Otherwise, keep original validated answer",
        ],
        9.2,
        4.95,
        3.5,
        1.65,
        (255, 240, 229),
    )
    add_source_ref(s, "utils/answer_generator.py")
    add_takeaway(s, "Answer quality is controlled by two gates: grounding gate and coverage-completeness gate.")

    # 7a. Answer flow zoom
    s = prs.slides.add_slide(blank)
    add_title(s, "6A. Answer Flow Zoom: Retry + Coverage Rewrite", "Readable view of regeneration and acceptance logic")
    add_image_crop(
        s,
        IMG["answer_detailed"],
        0.7,
        1.5,
        12.0,
        5.5,
        crop_left=0.18,
        crop_right=0.05,
        crop_top=0.08,
    )
    add_source_ref(s, "docs/qagredo_answer_generator_detailed.dot")
    add_takeaway(s, "Focus: grounding retries, coverage rewrite trigger, and revised-answer grounding gate.")

    # 8. Answer controls and trade-offs
    s = prs.slides.add_slide(blank)
    add_title(s, "7. Answer Controls and Trade-offs", "How answer quality behavior is tuned in config")
    add_card(
        s,
        "Grounding rejection loop",
        [
            "answer_generation.multi_turn.enable_rejection",
            "answer_generation.multi_turn.min_confidence_threshold",
            "answer_generation.multi_turn.max_regeneration_attempts",
            "Effect: stricter grounding -> higher confidence, higher latency",
        ],
        0.8,
        1.8,
        6.0,
        2.5,
        (232, 241, 255),
    )
    add_card(
        s,
        "Coverage rewrite loop",
        [
            "answer_generation.coverage_validation.enable",
            "answer_generation.coverage_validation.min_score_threshold",
            "answer_generation.coverage_validation.max_doc_chars",
            "Effect: better completeness, possible extra model call",
        ],
        0.8,
        4.45,
        6.0,
        2.1,
        (233, 252, 244),
    )
    add_card(
        s,
        "Practical interpretation",
        [
            "If many retries: reduce complexity or relax thresholds cautiously",
            "If many partial answers: increase coverage strictness",
            "If latency too high: lower retries before lowering quality gates",
        ],
        7.1,
        1.8,
        5.4,
        4.75,
        (245, 248, 255),
    )
    add_source_ref(s, "utils/answer_generator.py + config/config.yaml")
    add_takeaway(s, "Tune answer controls as a quality-latency dial, not as independent toggles.")

    # 9. Value proposition
    s = prs.slides.add_slide(blank)
    add_title(s, "8. Business Value", "Why teams adopt QAGRedo")
    add_bullets(
        s,
        [
            "Governance-ready: every decision can be traced to evidence and grading rationale.",
            "Offline-ready: suitable for air-gapped and restricted environments.",
            "Configurable rigor: tune for faster iteration or stricter quality controls.",
            "Operationally practical: run-level summaries simplify QA review and escalation.",
        ],
        x=0.8,
        y=1.8,
        w=6.2,
        h=4.6,
        font_size=17,
    )
    add_card(s, "Primary outcomes", ["Higher confidence in generated QA", "Better audit posture", "Clear failure visibility"], 7.3, 2.1, 5.2, 3.4, (245, 248, 255))
    add_takeaway(s, "The key value is controlled quality under operational constraints.")

    # 10. Tradeoff matrix
    s = prs.slides.add_slide(blank)
    add_title(s, "9. Trade-offs and Decision Guidance", "Choose settings by risk tolerance and latency budget")
    add_tradeoff_table(s)
    add_takeaway(s, "QAGRedo is strongest when quality and auditability matter more than raw speed.")

    # 11. Where it fits / not fit
    s = prs.slides.add_slide(blank)
    add_title(s, "10. Best-Fit vs Not-Ideal Scenarios", "Set expectations before rollout")
    add_card(
        s,
        "Best fit",
        [
            "Compliance and policy workflows",
            "Internal knowledge QA in restricted networks",
            "High-stakes reports requiring audit trails",
            "Teams willing to tune thresholds deliberately",
        ],
        0.8,
        1.8,
        5.9,
        4.6,
        (233, 252, 244),
    )
    add_card(
        s,
        "Not ideal (without simplification)",
        [
            "Ultra-low-latency interactive assistants",
            "Use cases needing open-web retrieval",
            "Teams without ops ownership for model services",
            "Situations prioritizing speed over evidential quality",
        ],
        7.0,
        1.8,
        5.4,
        4.6,
        (255, 240, 229),
    )
    add_takeaway(s, "Position QAGRedo as a quality-first pipeline, not a generic chatbot replacement.")

    # 12. Risks and mitigations
    s = prs.slides.add_slide(blank)
    add_title(s, "11. Key Risks and Mitigations", "Professional risk framing")
    add_card(s, "Risk: Runtime complexity", ["Mitigation: standard runbooks", "Mitigation: preset config profiles"], 0.8, 1.9, 4.0, 2.1, (246, 250, 255))
    add_card(s, "Risk: Latency increase", ["Mitigation: strict judge preflight", "Mitigation: targeted retries only"], 5.0, 1.9, 4.0, 2.1, (246, 250, 255))
    add_card(s, "Risk: Model mismatch/config drift", ["Mitigation: health/model checks in startup", "Mitigation: env + config consistency checks"], 9.2, 1.9, 3.3, 2.1, (246, 250, 255))
    add_image(s, IMG["part2"], 0.8, 4.2, 12.0, 2.1)
    add_takeaway(s, "Most risks are operational and manageable with clear run discipline.")

    # 13. KPI framework
    s = prs.slides.add_slide(blank)
    add_title(s, "12. KPI Framework for Leadership", "How to measure whether deployment is working")
    add_card(s, "Quality KPIs", ["Grounded vs ungrounded rate", "Grade distribution trend", "Coverage rewrite acceptance"], 0.8, 1.9, 3.9, 3.8, (232, 241, 255))
    add_card(s, "Efficiency KPIs", ["Total runtime per document", "Retry counters by stage", "Judge fallback frequency"], 4.95, 1.9, 3.9, 3.8, (233, 252, 244))
    add_card(s, "Governance KPIs", ["Audit completeness", "Failure reason completeness", "Reproducibility across reruns"], 9.1, 1.9, 3.3, 3.8, (240, 236, 255))
    add_takeaway(s, "Track quality and throughput together; optimize one without the other is misleading.")

    # 14. Rollout strategy
    s = prs.slides.add_slide(blank)
    add_title(s, "13. Rollout Plan (Recommended)", "Phase-based adoption to control risk")
    add_bullets(
        s,
        [
            "Phase 1: Baseline with strict llm-judge settings and representative documents.",
            "Phase 2: Tune question complexity and coverage thresholds by domain.",
            "Phase 3: Institutionalize KPI review and incident triage workflow.",
            "Phase 4: Standardize profiles (fast / standard / strict) by use case.",
        ],
        x=0.8,
        y=1.8,
        w=7.0,
        h=4.8,
        font_size=17,
    )
    add_card(s, "Change-management principle", ["Adjust one lever at a time", "Compare metrics before/after", "Avoid multi-variable drift"], 8.2, 2.2, 4.3, 3.0, (245, 248, 255))
    add_takeaway(s, "Controlled rollout gives predictable gains with low operational surprise.")

    # 15. Closing
    s = prs.slides.add_slide(blank)
    add_title(s, "14. Final Recommendation", "Use QAGRedo when trustworthiness is a first-class requirement")
    add_bullets(
        s,
        [
            "Adopt QAGRedo for quality-critical, auditable document QA workflows.",
            "Treat latency and compute as tunable costs tied to confidence requirements.",
            "Keep strict llm judging and independent judge model as default governance posture.",
            "Present trade-offs explicitly to stakeholders: quality, latency, cost, and operational complexity.",
        ],
        x=0.9,
        y=1.9,
        w=11.8,
        h=4.5,
        font_size=20,
    )
    final_note = s.shapes.add_textbox(Inches(0.9), Inches(6.0), Inches(11.7), Inches(0.8)).text_frame
    final_note.clear()
    p = final_note.paragraphs[0]
    p.text = "Message for leadership: QAGRedo is a quality-governance system with explicit trade-offs, not a black-box generator."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(22, 56, 108)
    p.alignment = PP_ALIGN.CENTER

    OUTFILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTFILE))
    print(str(OUTFILE))


if __name__ == "__main__":
    build()

