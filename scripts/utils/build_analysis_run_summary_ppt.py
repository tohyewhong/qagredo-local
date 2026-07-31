#!/usr/bin/env python3
"""Build a presentation summarizing one QAG doc_*_analysis.json file."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(15, 44, 96)
NAVY_BRIGHT = RGBColor(33, 83, 160)
SLATE = RGBColor(35, 48, 69)
MUTED = RGBColor(90, 104, 128)
PANEL = RGBColor(246, 250, 255)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(200, 214, 235)


def _trunc(text: str, max_len: int) -> str:
    t = (text or "").replace("\n", " ").strip()
    if len(t) <= max_len:
        return t
    return t[: max_len - 1] + "…"


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide) -> None:
    r = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    r.fill.solid()
    r.fill.fore_color.rgb = WHITE
    r.line.fill.background()


def _bar(slide) -> None:
    b = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(0.2),
        Inches(7.5),
    )
    b.fill.solid()
    b.fill.fore_color.rgb = NAVY_BRIGHT
    b.line.fill.background()


def _heading(slide, title: str, sub: str = "") -> None:
    t = slide.shapes.add_textbox(
        Inches(0.55), Inches(0.38), Inches(12.2), Inches(0.72)
    ).text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if sub:
        s = slide.shapes.add_textbox(
            Inches(0.55), Inches(1.0), Inches(12.0), Inches(0.4)
        ).text_frame
        s.clear()
        sp = s.paragraphs[0]
        sp.text = sub
        sp.font.size = Pt(12)
        sp.font.color.rgb = MUTED


def _footer(slide, text: str) -> None:
    f = slide.shapes.add_textbox(
        Inches(0.55), Inches(7.05), Inches(12.1), Inches(0.35)
    ).text_frame
    f.clear()
    p = f.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def _bullets(slide, items: List[str], x: float, y: float, w: float, h: float):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(13)
        para.font.color.rgb = SLATE
        para.space_after = Pt(6)
        para.line_spacing = 1.2


def _load(path: Path) -> Dict[str, Any]:
    with open(path, encoding="utf-8") as f:
        data = json.load(f)
    if not isinstance(data, dict):
        raise ValueError("Expected JSON object")
    return data


def build_prs(data: Dict[str, Any], src_path: Path) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    raw_doc = data.get("document")
    doc = raw_doc if isinstance(raw_doc, dict) else {}
    doc_id = str(doc.get("id", "—"))
    title = str(doc.get("title") or doc_id)
    gs = data.get("grading_summary") or {}
    grade = gs.get("overall_grade", "—")
    conf = gs.get("overall_confidence", "—")
    method = gs.get("grading_method", "—")
    judge = gs.get("judge_model") or "—"
    rm = data.get("run_metrics") or {}
    timings = (rm.get("timings_seconds") or {}) if isinstance(rm, dict) else {}
    qg = data.get("question_generation") or {}
    ag = data.get("answer_generation") or {}
    qp = data.get("qa_pairs")
    pairs = qp if isinstance(qp, list) else []

    # 1 — Title
    s = _blank(prs)
    _bg(s)
    _bar(s)
    tf = s.shapes.add_textbox(
        Inches(0.55), Inches(1.2), Inches(11.8), Inches(1.0)
    ).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "QAG run summary"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p = tf.add_paragraph()
    p.text = f"{doc_id} — {title}"
    p.font.size = Pt(20)
    p.font.color.rgb = SLATE
    p = tf.add_paragraph()
    p.text = f"Overall: grade {grade}, confidence {conf}, {method}"
    p.font.size = Pt(16)
    p.font.color.rgb = NAVY_BRIGHT
    p.font.bold = True
    _footer(s, _trunc(str(src_path), 95))

    # 2 — Output fields (schema)
    s = _blank(prs)
    _bg(s)
    _heading(
        s,
        "What’s in the JSON (top-level keys)",
        "One file per document after the full pipeline.",
    )
    keys = [
        "document — id, title, source, type, content used for Q&A",
        "qa_pairs[] — question, answer, hallucination_check, "
        "citation_spans, provenance",
        "question_generation — model, types, complexity, validation",
        "answer_generation — model, provider, timestamps, num_answers",
        "grading_summary — letter grade, confidence, method, judge",
        "run_metrics — stage timings (seconds), quality_counters",
        "Legacy: top-level hallucination_checks[] or pair.grading",
    ]
    _bullets(s, keys, 0.65, 1.45, 12.0, 5.5)
    _footer(s, "Schema: run_qa_pipeline combined_result")

    # 3 — Document excerpt
    s = _blank(prs)
    _bg(s)
    _heading(s, "Source document snapshot", doc_id)
    panel = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.52),
        Inches(1.35),
        Inches(12.2),
        Inches(5.55),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = LINE
    meta = [
        f"id: {doc_id}",
        f"title: {title}",
        f"source: {doc.get('source')!s}  type: {doc.get('type')!s}",
    ]
    content = doc.get("content") or ""
    excerpt = _trunc(str(content), 900)
    body = "\n".join(meta) + "\n\ncontent (excerpt):\n" + excerpt
    tb = s.shapes.add_textbox(
        Inches(0.72), Inches(1.52), Inches(11.85), Inches(5.25)
    ).text_frame
    tb.clear()
    tb.word_wrap = True
    tb.vertical_anchor = MSO_ANCHOR.TOP
    para = tb.paragraphs[0]
    para.text = body
    para.font.size = Pt(11)
    para.font.color.rgb = SLATE
    para.line_spacing = 1.15
    _footer(s, "document.content is what grading and citations use")

    # 4 — Grading + timing
    s = _blank(prs)
    _bg(s)
    _heading(s, "Grading summary & run metrics", "Rollups for this document")
    rows = 8
    cols = 2
    tbl = s.shapes.add_table(
        rows, cols, Inches(0.6), Inches(1.4), Inches(12.0), Inches(4.2)
    ).table
    tbl.columns[0].width = Inches(3.4)
    tbl.columns[1].width = Inches(8.6)
    hdr = ("Measure", "Value")
    cells = [
        ("overall_grade", str(grade)),
        ("overall_confidence", str(conf)),
        ("grading_method", str(method)),
        ("judge_model", str(judge)),
        ("time_question_gen_s", str(timings.get("question_generation", "—"))),
        ("time_answer_gen_s", str(timings.get("answer_generation", "—"))),
        ("time_grading_s", str(timings.get("grading", "—"))),
    ]
    for c, h in enumerate(hdr):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_BRIGHT
        cp = cell.text_frame.paragraphs[0]
        cp.font.bold = True
        cp.font.size = Pt(12)
        cp.font.color.rgb = WHITE
    for r, (k, v) in enumerate(cells, start=1):
        for c, val in enumerate((k, v)):
            cell = tbl.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if r % 2 else WHITE
            cp = cell.text_frame.paragraphs[0]
            cp.font.size = Pt(11)
            cp.font.color.rgb = SLATE
    qtypes = qg.get("question_types")
    qtypes_s = ", ".join(qtypes) if isinstance(qtypes, list) else str(qtypes)
    extra = (
        f"Questions: complexity={qg.get('complexity', '—')}, "
        f"types={_trunc(qtypes_s, 80)}"
    )
    _bullets(s, [extra], 0.6, 5.75, 12.0, 1.0)
    _footer(s, "answer_generation.num_answers = len(qa_pairs)")

    # 5 — QA table
    s = _blank(prs)
    _bg(s)
    _heading(
        s,
        f"QA pairs ({len(pairs)} slots)",
        "Question / answer excerpts and per-slot grounding.",
    )
    n = min(len(pairs), 5) + 1
    tbl2 = s.shapes.add_table(
        n, 4, Inches(0.5), Inches(1.35), Inches(12.25), Inches(5.45)
    ).table
    tbl2.columns[0].width = Inches(0.55)
    tbl2.columns[1].width = Inches(4.35)
    tbl2.columns[2].width = Inches(4.35)
    tbl2.columns[3].width = Inches(2.85)
    heads = ("#", "Question (excerpt)", "Answer (excerpt)", "Check")
    for c, h in enumerate(heads):
        cell = tbl2.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_BRIGHT
        cp = cell.text_frame.paragraphs[0]
        cp.font.bold = True
        cp.font.size = Pt(10)
        cp.font.color.rgb = WHITE
    for i in range(min(len(pairs), 5)):
        pair = pairs[i]
        gr = pair.get("hallucination_check")
        if not isinstance(gr, dict):
            gr = pair.get("grading")
        g = gr if isinstance(gr, dict) else {}
        meth = _trunc(str(g.get("method", "")), 28)
        grad_line = (
            f"g={g.get('is_grounded')} c={g.get('confidence')}\n{meth}"
        )
        row_vals = (
            str(i + 1),
            _trunc(str(pair.get("question", "")), 200),
            _trunc(str(pair.get("answer", "")), 200),
            grad_line,
        )
        for c, val in enumerate(row_vals):
            cell = tbl2.cell(i + 1, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if (i + 1) % 2 else WHITE
            cp = cell.text_frame.paragraphs[0]
            cp.font.size = Pt(9)
            cp.font.color.rgb = SLATE
    _footer(s, "Full text in JSON qa_pairs[].question / .answer")

    # 6 — Generation metadata
    s = _blank(prs)
    _bg(s)
    _heading(s, "Generation metadata", "Provenance for questions and answers")
    lines = [
        f"Question model: {qg.get('model', '—')} ({qg.get('provider', '—')})",
        f"Question time: {qg.get('timestamp', '—')} {qg.get('timezone', '')}",
        f"Answer model: {ag.get('model', '—')} ({ag.get('provider', '—')})",
        f"Answer time: {ag.get('timestamp', '—')} {ag.get('timezone', '')}",
        f"num_answers: {ag.get('num_answers', '—')}",
    ]
    chk = data.get("hallucination_checks")
    n_legacy = len(chk) if isinstance(chk, list) else 0
    n_slot = sum(
        1
        for p in pairs
        if isinstance(p, dict)
        and isinstance(p.get("hallucination_check"), dict)
    )
    lines.append(
        f"hallucination_check dicts in qa_pairs: {n_slot}; "
        f"legacy top-level list length: {n_legacy}"
    )
    _bullets(s, lines, 0.65, 1.45, 11.9, 5.2)
    _footer(s, str(src_path.name))

    return prs


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Build PPTX summary from one doc_*_analysis.json.",
    )
    parser.add_argument(
        "analysis",
        type=Path,
        help="Path to doc_*_analysis.json",
    )
    parser.add_argument(
        "--out",
        type=Path,
        default=None,
        help=(
            "Output .pptx (default: docs/QAG_Run_Summary_<doc_id>.pptx)"
        ),
    )
    args = parser.parse_args()
    path = args.analysis.expanduser().resolve()
    if not path.is_file():
        raise SystemExit(f"Not found: {path}")
    data = _load(path)
    raw_doc = data.get("document")
    doc = raw_doc if isinstance(raw_doc, dict) else {}
    stem = str(doc.get("id", "document")).replace("/", "_").replace(" ", "_")
    root = Path(__file__).resolve().parents[2]
    out = args.out
    if out is None:
        out = root / "docs" / f"QAG_Run_Summary_{stem}.pptx"
    else:
        out = out.expanduser().resolve()
    prs = build_prs(data, path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
